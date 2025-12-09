"""Stitch generated slide images into PowerPoint or PDF."""

import logging
from pathlib import Path
from typing import List, Dict, Optional
from datetime import datetime

logger = logging.getLogger(__name__)


class SlideStitcher:
    """Stitch individual slide images into a complete presentation."""
    
    def __init__(self):
        self.output_dir = Path("exports")
        self.output_dir.mkdir(exist_ok=True)
    
    def stitch_images_to_pptx(self, slide_images: Dict[int, str], company_name: str) -> Optional[str]:
        """
        Stitch slide images into a PowerPoint presentation.
        
        Args:
            slide_images: Dictionary mapping slide_number to image_path
            company_name: Company name for filename
            
        Returns:
            Path to created PowerPoint file, or None if failed
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            logger.info(f"📎 Stitching {len(slide_images)} slide images into PowerPoint...")
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Sort slides by slide number
            sorted_slides = sorted(slide_images.items(), key=lambda x: x[0])
            
            for slide_num, image_path in sorted_slides:
                # Handle BytesIO objects - convert to file path if needed
                if hasattr(image_path, 'read'):
                    # It's a BytesIO object, need to save it first
                    from io import BytesIO
                    import tempfile
                    # Reset file pointer to beginning
                    image_path.seek(0)
                    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                    temp_file.write(image_path.read())
                    temp_file.close()
                    image_path = temp_file.name
                    logger.info(f"💾 Saved BytesIO image to temporary file: {image_path}")
                # Also check if it's a Path object
                elif isinstance(image_path, Path):
                    image_path = str(image_path)
                
                # Ensure it's a string path
                image_path = str(image_path)
                
                if not image_path or not Path(image_path).exists():
                    logger.warning(f"⚠️ Image not found for slide {slide_num}: {image_path}")
                    continue
                
                try:
                    # Create blank slide
                    blank_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # Add image to fill entire slide
                    left = top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height
                    
                    slide.shapes.add_picture(image_path, left, top, width, height)
                    logger.info(f"✅ Added slide {slide_num} image to PowerPoint: {image_path}")
                    
                except Exception as e:
                    logger.error(f"Error adding slide {slide_num} image: {e}")
                    continue
            
            # Save PowerPoint
            filename = f"{company_name.replace(' ', '_')}_stitched_slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = self.output_dir / filename
            
            prs.save(str(filepath))
            logger.info(f"✅ Stitched PowerPoint saved: {filepath}")
            return str(filepath)
            
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"Error stitching images to PowerPoint: {e}", exc_info=True)
            return None
    
    def stitch_images_to_pdf(self, slide_images: Dict[int, str], company_name: str) -> Optional[str]:
        """
        Stitch slide images into a PDF presentation.
        
        Args:
            slide_images: Dictionary mapping slide_number to image_path
            company_name: Company name for filename
            
        Returns:
            Path to created PDF file, or None if failed
        """
        try:
            from PIL import Image
            
            logger.info(f"📎 Stitching {len(slide_images)} slide images into PDF...")
            
            # Sort slides by slide number
            sorted_slides = sorted(slide_images.items(), key=lambda x: x[0])
            
            # Open all images
            images = []
            for slide_num, image_path in sorted_slides:
                if not image_path or not Path(image_path).exists():
                    logger.warning(f"⚠️ Image not found for slide {slide_num}: {image_path}")
                    continue
                
                try:
                    img = Image.open(image_path)
                    # Convert to RGB if needed (for PDF)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
                    logger.info(f"✅ Loaded slide {slide_num} image")
                except Exception as e:
                    logger.error(f"Error loading slide {slide_num} image: {e}")
                    continue
            
            if not images:
                logger.error("No valid images to stitch")
                return None
            
            # Save as PDF
            filename = f"{company_name.replace(' ', '_')}_stitched_slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            filepath = self.output_dir / filename
            
            images[0].save(
                str(filepath),
                "PDF",
                resolution=100.0,
                save_all=True,
                append_images=images[1:] if len(images) > 1 else []
            )
            
            logger.info(f"✅ Stitched PDF saved: {filepath}")
            return str(filepath)
            
        except ImportError:
            logger.error("Pillow (PIL) not installed. Install with: pip install Pillow")
            return None
        except Exception as e:
            logger.error(f"Error stitching images to PDF: {e}", exc_info=True)
            return None

