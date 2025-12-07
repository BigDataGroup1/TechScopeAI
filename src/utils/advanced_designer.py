"""Advanced visual design elements for industry-level PowerPoint presentations."""

import logging
from typing import Dict, Tuple, Optional
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

logger = logging.getLogger(__name__)


class AdvancedDesigner:
    """Advanced visual design elements for professional presentations."""
    
    def __init__(self):
        """Initialize advanced designer."""
        logger.info("AdvancedDesigner initialized")
    
    def create_gradient_fill(self, shape, color1: RGBColor, color2: RGBColor,
                           direction: str = "vertical") -> None:
        """
        Create gradient fill for a shape.
        
        Args:
            shape: PowerPoint shape object
            color1: Start color
            color2: End color
            direction: "vertical" or "horizontal"
        """
        try:
            fill = shape.fill
            fill.gradient()
            fill.gradient_angle = 90.0 if direction == "vertical" else 0.0
            
            # Set gradient stops
            fill.gradient_stops[0].color.rgb = color1
            fill.gradient_stops[1].color.rgb = color2
        except Exception as e:
            logger.warning(f"Could not create gradient, using solid fill: {e}")
            shape.fill.solid()
            shape.fill.fore_color.rgb = color1
    
    def add_shadow(self, shape, blur_radius: int = 5, 
                   offset_x: int = 2, offset_y: int = 2,
                   color: Optional[RGBColor] = None) -> None:
        """
        Add shadow effect to shape.
        
        Args:
            shape: PowerPoint shape object
            blur_radius: Shadow blur amount
            offset_x: Horizontal offset
            offset_y: Vertical offset
            color: Shadow color (default: black with transparency)
        """
        try:
            shadow = shape.shadow
            shadow.inherit = False
            shadow.style = 'outer'
            shadow.blur_radius = Pt(blur_radius)
            shadow.offset_x = Pt(offset_x)
            shadow.offset_y = Pt(offset_y)
            if color:
                shadow.color.rgb = color
            else:
                # Default: black with some transparency
                shadow.color.rgb = RGBColor(0, 0, 0)
        except Exception as e:
            logger.warning(f"Could not add shadow: {e}")
    
    def create_rounded_rectangle(self, slide, left, top, width, height,
                                corner_radius: float = 0.1,
                                fill_color: Optional[RGBColor] = None) -> object:
        """
        Create rounded rectangle shape.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            corner_radius: Corner radius (0.0 to 1.0)
            fill_color: Fill color
            
        Returns:
            Shape object
        """
        # python-pptx doesn't support rounded rectangles directly
        # Use regular rectangle with adjusted line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        
        # Adjust line for rounded effect (visual only)
        shape.line.color.rgb = fill_color if fill_color else RGBColor(200, 200, 200)
        shape.line.width = Pt(1)
        
        return shape
    
    def add_icon_placeholder(self, slide, left, top, size: float = 1.0,
                           icon_type: str = "check") -> object:
        """
        Add icon placeholder (using shapes as icons).
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            size: Icon size in inches
            icon_type: Type of icon (check, star, arrow, etc.)
            
        Returns:
            Shape object representing icon
        """
        icon_shapes = {
            "check": MSO_SHAPE.ROUNDED_RECTANGLE,  # Placeholder
            "star": MSO_SHAPE.STAR_5_POINT,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "circle": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.TRIANGLE,
            "diamond": MSO_SHAPE.DIAMOND
        }
        
        shape_type = icon_shapes.get(icon_type, MSO_SHAPE.OVAL)
        shape = slide.shapes.add_shape(shape_type, left, top, Inches(size), Inches(size))
        
        # Style the icon
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(59, 130, 246)  # Blue
        shape.line.fill.background()
        
        return shape
    
    def create_text_with_background(self, slide, left, top, width, height,
                                   text: str, bg_color: RGBColor,
                                   text_color: RGBColor = RGBColor(255, 255, 255),
                                   padding: float = 0.2) -> object:
        """
        Create text box with colored background.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Text content
            bg_color: Background color
            text_color: Text color
            padding: Padding in inches
            
        Returns:
            Text box shape
        """
        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        
        # Text box on top
        text_box = slide.shapes.add_textbox(
            left + Inches(padding),
            top + Inches(padding),
            width - Inches(padding * 2),
            height - Inches(padding * 2)
        )
        text_frame = text_box.text_frame
        text_frame.text = text
        para = text_frame.paragraphs[0]
        para.font.color.rgb = text_color
        para.font.size = Pt(14)
        para.font.bold = True
        
        return text_box
    
    def create_visual_separator(self, slide, left, top, width,
                              color: RGBColor = RGBColor(200, 200, 200),
                              thickness: float = 0.05) -> object:
        """
        Create visual separator line.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            color: Line color
            thickness: Line thickness in inches
            
        Returns:
            Shape object
        """
        separator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, Inches(thickness)
        )
        separator.fill.solid()
        separator.fill.fore_color.rgb = color
        separator.line.fill.background()
        
        return separator
    
    def create_stat_card(self, slide, left, top, width, height,
                        value: str, label: str,
                        bg_color: RGBColor = RGBColor(59, 130, 246),
                        icon_type: str = "circle") -> object:
        """
        Create a stat card with value and label.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            value: Stat value (e.g., "10M", "45%")
            label: Stat label (e.g., "Users", "Growth")
            bg_color: Background color
            icon_type: Icon type
            
        Returns:
            Shape object
        """
        # Background with gradient
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        self.create_gradient_fill(card, bg_color, self._darken_color(bg_color))
        self.add_shadow(card)
        
        # Icon
        icon_size = 0.5
        icon = self.add_icon_placeholder(
            slide, left + Inches(0.2), top + Inches(0.2), icon_size, icon_type
        )
        
        # Value text
        value_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), Inches(0.8)
        )
        value_frame = value_box.text_frame
        value_frame.text = value
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(32)
        value_para.font.bold = True
        value_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Label text
        label_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(1.6), width - Inches(0.4), Inches(0.6)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = RGBColor(255, 255, 255)
        
        return card
    
    def _darken_color(self, color: RGBColor, factor: float = 0.8) -> RGBColor:
        """Darken a color by a factor."""
        return RGBColor(
            int(color.r * factor),
            int(color.g * factor),
            int(color.b * factor)
        )
    
    def create_progress_bar(self, slide, left, top, width, height,
                           percentage: float,
                           fill_color: RGBColor = RGBColor(34, 197, 94),
                           bg_color: RGBColor = RGBColor(229, 231, 235)) -> object:
        """
        Create a progress bar.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Total width in inches
            height: Height in inches
            percentage: Progress percentage (0-100)
            fill_color: Fill color
            bg_color: Background color
            
        Returns:
            Progress bar shape
        """
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        
        # Fill
        fill_width = width * (percentage / 100.0)
        fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, fill_width, height)
        fill.fill.solid()
        fill.fill.fore_color.rgb = fill_color
        fill.line.fill.background()
        
        # Percentage text
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = f"{percentage:.0f}%"
        para = text_frame.paragraphs[0]
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = RGBColor(0, 0, 0)
        para.alignment = 1  # Center
        
        return fill
    
    def create_badge(self, slide, left, top, text: str,
                    bg_color: RGBColor = RGBColor(34, 197, 94),
                    text_color: RGBColor = RGBColor(255, 255, 255)) -> object:
        """
        Create a badge/pill shape.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            text: Badge text
            bg_color: Background color
            text_color: Text color
            
        Returns:
            Badge shape
        """
        # Estimate width based on text length
        width = Inches(0.3 * len(text) + 0.4)
        height = Inches(0.4)
        
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        badge.fill.solid()
        badge.fill.fore_color.rgb = bg_color
        badge.line.fill.background()
        
        # Text
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = text
        para = text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = text_color
        para.alignment = 1  # Center
        
        return badge


