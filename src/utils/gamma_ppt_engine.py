"""Gamma.ai-style PowerPoint engine with smart layouts, charts, and templates."""

import logging
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from .layout_engine import LayoutEngine
from .template_manager import TemplateManager
from .chart_generator import ChartGenerator
from .advanced_designer import AdvancedDesigner
from .content_enhancer import ContentEnhancer
from .brand_manager import BrandManager
from .image_enhancer import ImageEnhancer
from .infographic_generator import InfographicGenerator

logger = logging.getLogger(__name__)


class GammaPPTEngine:
    """Gamma.ai-style PowerPoint generation engine with industry-level features."""
    
    def __init__(self, template_name: str = "modern",
                 logo_path: Optional[str] = None,
                 brand_colors: Optional[Dict[str, str]] = None):
        """
        Initialize Gamma PPT engine.
        
        Args:
            template_name: Template to use (modern, bold, corporate, startup)
            logo_path: Optional path to company logo
            brand_colors: Optional brand colors dict (primary, secondary)
        """
        self.layout_engine = LayoutEngine()
        self.template_manager = TemplateManager(template_name)
        self.chart_generator = ChartGenerator()
        self.designer = AdvancedDesigner()
        self.content_enhancer = ContentEnhancer()
        self.image_enhancer = ImageEnhancer()
        self.infographic_generator = InfographicGenerator()
        
        # Brand customization
        primary_color = brand_colors.get('primary') if brand_colors else None
        secondary_color = brand_colors.get('secondary') if brand_colors else None
        self.brand_manager = BrandManager(
            logo_path=logo_path,
            primary_color=primary_color,
            secondary_color=secondary_color
        )
        
        logger.info(f"GammaPPTEngine initialized with template: {template_name}")
    
    def create_presentation(self, company_name: str, slides: List[Dict],
                          include_images: bool = True,
                          include_charts: bool = True) -> Presentation:
        """
        Create a Gamma.ai-style presentation with industry-level features.
        
        Args:
            company_name: Company name
            slides: List of slide dictionaries
            include_images: Whether to include images
            include_charts: Whether to auto-generate charts
            
        Returns:
            PowerPoint Presentation object
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Get template colors (with brand customization)
        default_colors = self.template_manager.get_colors()
        colors = self.brand_manager.get_brand_colors(default_colors)
        fonts = self.brand_manager.get_brand_fonts(self.template_manager.get_fonts())
        
        # Create title slide with advanced design
        self._create_enhanced_title_slide(prs, company_name, colors, fonts)
        
        # Create content slides with all enhancements
        for i, slide_data in enumerate(slides):
            layout_config = self.layout_engine.get_layout_for_slide(
                slide_data.get('title', ''),
                slide_data.get('content', ''),
                slide_data.get('key_points', [])
            )
            
            # Enhance content
            enhanced_content = self._enhance_slide_content(slide_data, company_name)
            slide_data['enhanced_content'] = enhanced_content
            
            self._create_enhanced_content_slide(
                prs, slide_data, layout_config, colors, fonts,
                include_images, include_charts, company_name
            )
        
        return prs
    
    def _enhance_slide_content(self, slide_data: Dict, company_name: str) -> Dict:
        """Enhance slide content with hooks and storytelling."""
        slide_title = slide_data.get('title', '')
        slide_content = slide_data.get('content', '')
        key_points = slide_data.get('key_points', [])
        
        # Generate hook
        hook = self.content_enhancer.generate_hook(slide_title, slide_content, company_name)
        
        # Enhance content with storytelling
        enhanced_content = self.content_enhancer.enhance_with_storytelling(
            slide_content, slide_title
        )
        
        # Extract metrics
        metrics = self.content_enhancer.extract_metrics(slide_content)
        
        # Enhance key points
        enhanced_key_points = self.content_enhancer.enhance_key_points(key_points)
        
        return {
            'hook': hook,
            'enhanced_content': enhanced_content,
            'metrics': metrics,
            'enhanced_key_points': enhanced_key_points
        }
    
    def _create_enhanced_title_slide(self, prs: Presentation, company_name: str,
                                    colors: Dict, fonts: Dict):
        """Create enhanced title slide with gradients and shadows."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Full background with gradient
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        self.designer.create_gradient_fill(
            background, colors["secondary"], colors["primary"], "vertical"
        )
        background.line.fill.background()
        
        # Accent bar at top with shadow
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = colors["accent"]
        accent_bar.line.fill.background()
        self.designer.add_shadow(accent_bar)
        
        # Company name - large and bold with shadow effect
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = company_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(68)  # Larger
        title_para.font.bold = True
        title_para.font.color.rgb = colors["white"]
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle with better styling
        top = Inches(4.2)
        height = Inches(0.8)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Pitch Deck"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(203, 213, 225)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add logo if available
        self.brand_manager.add_logo_to_slide(slide, position="top-right", size=1.0)
    
    def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict,
                                      layout_config: Dict, colors: Dict, fonts: Dict,
                                      include_images: bool, include_charts: bool,
                                      company_name: str):
        """Create enhanced content slide with all new features."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        layout_type = layout_config["layout_type"]
        config = layout_config["config"]
        
        slide_title = slide_data.get('title', 'Untitled')
        slide_content = slide_data.get('content', '')
        key_points = slide_data.get('key_points', [])
        enhanced = slide_data.get('enhanced_content', {})
        
        # Enhanced header with gradients
        self._add_enhanced_header(slide, slide_title, slide_data.get('slide_number', 0),
                                 colors, fonts, prs.slide_width)
        
        # Add logo to every slide
        self.brand_manager.add_logo_to_slide(slide, position="top-right", size=0.6)
        
        # Generate chart/infographic if needed
        chart_path = None
        if include_charts:
            if config.get('has_chart', False):
                chart_path = self.chart_generator.auto_generate_chart_for_slide(
                    slide_title, slide_content, key_points
                )
            elif layout_type == "timeline":
                # Generate timeline infographic
                events = self._extract_timeline_events(slide_content, key_points)
                if events:
                    chart_path = self.infographic_generator.generate_timeline(
                        events, slide_title
                    )
            elif layout_type == "comparison":
                # Generate comparison matrix
                comparison_data = self._extract_comparison_data(slide_content, key_points)
                if comparison_data:
                    chart_path = self.infographic_generator.generate_comparison_matrix(
                        comparison_data['items'],
                        comparison_data['features'],
                        comparison_data['matrix'],
                        slide_title
                    )
        
        # Add content based on layout with enhancements
        if layout_type == "data" and chart_path:
            self._add_enhanced_data_slide(slide, slide_title, slide_content, key_points,
                                        chart_path, colors, fonts, prs.slide_width, enhanced)
        elif layout_type == "traction":
            self._add_enhanced_traction_slide(slide, slide_title, slide_content, key_points,
                                             colors, fonts, prs.slide_width, enhanced)
        elif layout_type == "comparison" and chart_path:
            self._add_enhanced_comparison_slide(slide, slide_title, slide_content, key_points,
                                              chart_path, colors, fonts, prs.slide_width)
        else:
            # Enhanced default layout
            self._add_enhanced_default_slide(slide, slide_title, slide_content, key_points,
                                            slide_data.get('image_path'), colors, fonts,
                                            prs.slide_width, config, enhanced)
    
    def _add_enhanced_header(self, slide, title: str, slide_num: int,
                           colors: Dict, fonts: Dict, slide_width):
        """Add enhanced header with gradients and shadows."""
        # Header background with gradient
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                       slide_width, Inches(1.3))
        self.designer.create_gradient_fill(
            header, colors["primary"], colors["secondary"], "horizontal"
        )
        header.line.fill.background()
        self.designer.add_shadow(header, blur_radius=8, offset_x=0, offset_y=2)
        
        # Accent line with shadow
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                       slide_width, Inches(0.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = colors["accent"]
        accent.line.fill.background()
        
        # Slide number badge with rounded corners and shadow
        if slide_num > 0:
            num_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(0.3), Inches(0.4),
                                          Inches(0.9), Inches(0.5))
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = colors["white"]
            num_bg.line.fill.background()
            self.designer.add_shadow(num_bg)
            
            num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.4), Inches(0.9), Inches(0.5))
            num_frame = num_box.text_frame
            num_frame.text = str(slide_num)
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(24)
            num_para.font.bold = True
            num_para.font.color.rgb = colors["primary"]
            num_para.alignment = PP_ALIGN.CENTER
        
        # Title with better typography
        title_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.35), Inches(7.4), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(fonts["size_heading"] + 2)  # Slightly larger
        title_para.font.bold = True
        title_para.font.color.rgb = colors["white"]
    
    def _add_enhanced_default_slide(self, slide, title: str, content: str, key_points: List[str],
                                  image_path: Optional[str], colors: Dict, fonts: Dict,
                                  slide_width, config: Dict, enhanced: Dict):
        """Add enhanced default layout slide."""
        # Add hook if available
        hook = enhanced.get('hook', '')
        metrics = enhanced.get('metrics', {})
        enhanced_key_points = enhanced.get('enhanced_key_points', key_points)
        
        # Content area with better spacing
        left_content = Inches(0.6)
        top_content = Inches(1.7)
        width_content = Inches(4.5) if image_path else Inches(8.8)
        height_content = Inches(5.3)
        
        content_box = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add hook
        if hook:
            p = content_frame.paragraphs[0]
            p.text = hook
            p.font.size = Pt(fonts["size_body"] + 2)
            p.font.bold = True
            p.font.color.rgb = colors["accent"]
            p.space_after = Pt(12)
        
        # Add main content
        if content:
            if not hook:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = content
            p.font.size = Pt(fonts["size_body"])
            p.font.color.rgb = colors["text"]
            p.space_after = Pt(16)
            p.line_spacing = 1.3
        
        # Add metrics as stat cards if available
        if metrics and len(metrics) <= 3:
            metrics_y = top_content + Inches(3.5)
            metric_width = Inches(2.5)
            for i, (metric_name, value) in enumerate(list(metrics.items())[:3]):
                metric_x = left_content + (i * Inches(2.7))
                # Create stat card
                self.designer.create_stat_card(
                    slide, metric_x, metrics_y, metric_width, Inches(1.2),
                    f"{value:.1f}", metric_name.replace('_', ' ').title(),
                    colors["accent"]
                )
        else:
            # Add enhanced key points with icons
            start_idx = 1 if content else 0
            for i, point in enumerate(enhanced_key_points):
                if i == 0 and not content and not hook:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(fonts["size_body"] - 1)
                p.font.color.rgb = colors["text"]
                p.space_after = Pt(12)
                p.space_before = Pt(4)
                p.line_spacing = 1.4
        
        # Add enhanced image if available
        if image_path and Path(image_path).exists():
            try:
                # Enhance image (add overlay if needed)
                enhanced_image = self.image_enhancer.add_text_overlay(
                    image_path, title, position="bottom"
                ) or image_path
                
                left_img = Inches(5.3)
                top_img = Inches(1.7)
                width_img = Inches(4.4)
                height_img = Inches(5.5)
                
                # Add rounded corners effect
                img = slide.shapes.add_picture(
                    enhanced_image, left_img, top_img, width_img, height_img
                )
                
                # Add shadow effect (visual)
                shadow_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left_img + Inches(0.1),
                    top_img + Inches(0.1),
                    width_img, height_img
                )
                shadow_bg.fill.solid()
                shadow_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
                shadow_bg.line.fill.background()
                slide.shapes._spTree.remove(shadow_bg._element)
                slide.shapes._spTree.insert(2, shadow_bg._element)  # Send to back
                
            except Exception as e:
                logger.warning(f"Could not add enhanced image: {e}")
                # Fallback to regular image
                try:
                    slide.shapes.add_picture(
                        image_path, Inches(5.3), Inches(1.7),
                        Inches(4.4), Inches(5.5)
                    )
                except:
                    pass
    
    def _add_enhanced_data_slide(self, slide, title: str, content: str, key_points: List[str],
                                chart_path: Path, colors: Dict, fonts: Dict, slide_width,
                                enhanced: Dict):
        """Add enhanced data slide with chart and metrics."""
        # Chart in center with shadow
        if chart_path and Path(chart_path).exists():
            try:
                left_chart = Inches(0.8)
                top_chart = Inches(2)
                width_chart = Inches(8.4)
                height_chart = Inches(4.5)
                
                # Add shadow background
                shadow = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left_chart + Inches(0.05),
                    top_chart + Inches(0.05),
                    width_chart, height_chart
                )
                shadow.fill.solid()
                shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
                shadow.line.fill.background()
                slide.shapes._spTree.remove(shadow._element)
                slide.shapes._spTree.insert(2, shadow._element)
                
                slide.shapes.add_picture(
                    str(chart_path), left_chart, top_chart, width_chart, height_chart
                )
            except Exception as e:
                logger.warning(f"Could not add chart: {e}")
        
        # Add key metrics below chart
        metrics = enhanced.get('metrics', {})
        if metrics:
            metrics_y = Inches(6.5)
            metric_width = Inches(2.5)
            for i, (metric_name, value) in enumerate(list(metrics.items())[:4]):
                metric_x = Inches(0.8) + (i * Inches(2.3))
                self.designer.create_stat_card(
                    slide, metric_x, metrics_y, metric_width, Inches(0.8),
                    f"{value:.1f}", metric_name.replace('_', ' ').title()[:10],
                    colors["highlight"]
                )
    
    def _add_enhanced_traction_slide(self, slide, title: str, content: str, key_points: List[str],
                                    colors: Dict, fonts: Dict, slide_width, enhanced: Dict):
        """Add enhanced traction slide with stat cards."""
        metrics = enhanced.get('metrics', {})
        
        if metrics:
            # Create grid of stat cards
            num_metrics = len(metrics)
            cols = min(3, num_metrics)
            rows = (num_metrics + cols - 1) // cols
            
            card_width = Inches(2.8)
            card_height = Inches(1.8)
            spacing = Inches(0.3)
            
            start_x = (slide_width - (cols * card_width + (cols - 1) * spacing)) / 2
            start_y = Inches(2.5)
            
            for i, (metric_name, value) in enumerate(metrics.items()):
                row = i // cols
                col = i % cols
                
                x = start_x + col * (card_width + spacing)
                y = start_y + row * (card_height + spacing)
                
                self.designer.create_stat_card(
                    slide, x, y, card_width, card_height,
                    f"{value:.1f}", metric_name.replace('_', ' ').title(),
                    colors["accent"]
                )
        else:
            # Fallback to default layout
            self._add_enhanced_default_slide(
                slide, title, content, key_points, None, colors, fonts, slide_width, {}, enhanced
            )
    
    def _add_enhanced_comparison_slide(self, slide, title: str, content: str, key_points: List[str],
                                      chart_path: Path, colors: Dict, fonts: Dict, slide_width):
        """Add enhanced comparison slide."""
        # Comparison chart/image in center
        if chart_path and Path(chart_path).exists():
            try:
                left_chart = Inches(1)
                top_chart = Inches(2)
                width_chart = Inches(8)
                height_chart = Inches(5)
                slide.shapes.add_picture(
                    str(chart_path), left_chart, top_chart, width_chart, height_chart
                )
            except Exception as e:
                logger.warning(f"Could not add comparison chart: {e}")
    
    def _extract_timeline_events(self, content: str, key_points: List[str]) -> List[Dict]:
        """Extract timeline events from content."""
        # Simple extraction - look for dates and milestones
        import re
        events = []
        
        # Pattern for dates
        date_pattern = r'(\d{4}|\w+\s+\d{4}|Q\d\s+\d{4})'
        
        text = f"{content} {' '.join(key_points)}"
        sentences = re.split(r'[.!?]\s+', text)
        
        for sentence in sentences[:5]:  # Max 5 events
            dates = re.findall(date_pattern, sentence)
            if dates:
                events.append({
                    'date': dates[0],
                    'title': sentence[:50],
                    'description': sentence[50:100] if len(sentence) > 50 else ''
                })
        
        return events if events else None
    
    def _extract_comparison_data(self, content: str, key_points: List[str]) -> Optional[Dict]:
        """Extract comparison data from content."""
        # Simple extraction - look for vs/comparison patterns
        if 'vs' not in content.lower() and 'versus' not in content.lower():
            return None
        
        # Create simple comparison matrix
        items = ['Us', 'Competitor']
        features = key_points[:5] if key_points else ['Feature 1', 'Feature 2']
        matrix = [
            [True] * len(features),  # We have all features
            [False] * len(features)  # Competitor missing some
        ]
        
        return {
            'items': items,
            'features': features,
            'matrix': matrix
        }
    
    def save_presentation(self, prs: Presentation, company_name: str,
                         output_dir: Path = None) -> str:
        """
        Save presentation to file.
        
        Args:
            prs: Presentation object
            company_name: Company name for filename
            output_dir: Output directory (default: exports/)
            
        Returns:
            Path to saved file
        """
        if not output_dir:
            output_dir = Path("exports")
        output_dir.mkdir(exist_ok=True)
        
        filename = f"{company_name.replace(' ', '_')}_gamma_pitch_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = output_dir / filename
        
        prs.save(str(filepath))
        logger.info(f"Gamma-style presentation saved: {filepath}")
        return str(filepath)
