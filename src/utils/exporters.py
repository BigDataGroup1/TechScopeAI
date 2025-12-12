"""Export utilities for pitch decks."""

import logging
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime

logger = logging.getLogger(__name__)


class PitchExporter:
    """Export pitch decks to various formats."""
    
    def __init__(self):
        self.output_dir = Path("exports")
        self.output_dir.mkdir(exist_ok=True)
    
    def export_to_markdown(self, slides: List[Dict], company_name: str) -> str:
        """Export slides to Markdown format."""
        logger.info(f"Exporting {len(slides)} slides to Markdown")
        
        md_content = f"# Pitch Deck: {company_name}\n\n"
        md_content += f"*Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n"
        md_content += "---\n\n"
        
        for slide in slides:
            md_content += f"## Slide {slide.get('slide_number', '?')}: {slide.get('title', 'Untitled')}\n\n"
            
            content = slide.get('content', '')
            if content:
                md_content += f"{content}\n\n"
            
            key_points = slide.get('key_points', [])
            if key_points:
                md_content += "**Key Points:**\n"
                for point in key_points:
                    md_content += f"- {point}\n"
                md_content += "\n"
            
            md_content += "---\n\n"
        
        # Save to file
        filename = f"{company_name.replace(' ', '_')}_pitch_deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
        filepath = self.output_dir / filename
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(md_content)
        
        logger.info(f"Markdown exported to: {filepath}")
        return str(filepath)
    
    def export_to_pdf(self, slides: List[Dict], company_name: str) -> Optional[str]:
        """Export slides to PDF format."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.enums import TA_CENTER, TA_LEFT
            
            logger.info(f"Exporting {len(slides)} slides to PDF")
            
            filename = f"{company_name.replace(' ', '_')}_pitch_deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            filepath = self.output_dir / filename
            
            doc = SimpleDocTemplate(str(filepath), pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Title style
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor='#1a1a1a',
                spaceAfter=30,
                alignment=TA_CENTER
            )
            
            # Slide title style
            slide_title_style = ParagraphStyle(
                'SlideTitle',
                parent=styles['Heading2'],
                fontSize=18,
                textColor='#2c3e50',
                spaceAfter=12,
                alignment=TA_LEFT
            )
            
            # Content style
            content_style = ParagraphStyle(
                'Content',
                parent=styles['Normal'],
                fontSize=11,
                leading=14,
                spaceAfter=12,
                alignment=TA_LEFT
            )
            
            # Add title page
            story.append(Paragraph(f"<b>{company_name}</b>", title_style))
            story.append(Spacer(1, 0.5*inch))
            story.append(Paragraph("Pitch Deck", styles['Heading2']))
            story.append(Spacer(1, 0.3*inch))
            story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
            story.append(PageBreak())
            
            # Add slides
            for slide in slides:
                slide_num = slide.get('slide_number', '?')
                title = slide.get('title', 'Untitled')
                
                story.append(Paragraph(f"Slide {slide_num}: {title}", slide_title_style))
                story.append(Spacer(1, 0.2*inch))
                
                content = slide.get('content', '')
                if content:
                    # Replace newlines with HTML breaks
                    content = content.replace('\n', '<br/>')
                    story.append(Paragraph(content, content_style))
                    story.append(Spacer(1, 0.15*inch))
                
                key_points = slide.get('key_points', [])
                if key_points:
                    story.append(Paragraph("<b>Key Points:</b>", styles['Heading3']))
                    for point in key_points:
                        story.append(Paragraph(f"• {point}", content_style))
                    story.append(Spacer(1, 0.15*inch))
                
                story.append(PageBreak())
            
            doc.build(story)
            logger.info(f"PDF exported to: {filepath}")
            return str(filepath)
            
        except ImportError:
            logger.warning("reportlab not installed. Install with: pip install reportlab")
            return None
    
    def export_to_powerpoint(self, slides: List[Dict], company_name: str, 
                           include_images: bool = True, 
                           enhance_with_ai: bool = True,
                           company_data: Optional[Dict] = None,
                           full_rewrite: bool = True,
                           ai_provider: str = "auto") -> Optional[str]:
        """
        Export slides to PowerPoint format with professional styling.
        
        Args:
            slides: List of slide dictionaries
            company_name: Company name
            include_images: Whether to fetch and include images
            enhance_with_ai: Whether to enhance slides with AI before creating PPTX
            company_data: Optional company data for financial charts and intelligent filtering
            full_rewrite: Whether to do full AI rewrite pass (default: True)
            ai_provider: AI provider to use - "auto", "gemini", or "openai" (default: "auto")
            
        Returns:
            Path to created PowerPoint file
        """
        try:
            # 🆕 FULL REWRITE WITH AI (OpenAI or Gemini) - This is the main enhancement
            if enhance_with_ai and full_rewrite:
                enhancer = None
                provider_name = None
                
                # Determine which provider to use based on preference
                use_gemini = ai_provider.lower() in ["auto", "gemini"]
                use_openai = ai_provider.lower() in ["auto", "openai"]
                
                # Try Gemini first (if requested)
                if use_gemini:
                    try:
                        from .gemini_enhancer import GeminiSlideEnhancer
                        gemini_enhancer = GeminiSlideEnhancer()
                        if gemini_enhancer.enabled:
                            enhancer = gemini_enhancer
                            provider_name = "Gemini"
                            logger.info("🤖 Using Gemini for full pitch deck rewrite...")
                    except Exception as e:
                        logger.debug(f"Gemini not available: {e}")
                
                # Fallback to OpenAI if Gemini not available or OpenAI requested
                if not enhancer and use_openai:
                    try:
                        from .openai_enhancer import OpenAISlideEnhancer
                        openai_enhancer = OpenAISlideEnhancer()
                        if openai_enhancer.enabled:
                            enhancer = openai_enhancer
                            provider_name = "ChatGPT"
                            logger.info("🤖 Using ChatGPT for full pitch deck rewrite...")
                    except Exception as e:
                        logger.debug(f"OpenAI not available: {e}")
                
                # Perform rewrite if enhancer is available
                if enhancer:
                    try:
                        logger.info(f"🤖 Full {provider_name} rewrite of pitch deck...")
                        slides = enhancer.rewrite_full_pitch_deck(slides, company_name, company_data)
                        logger.info(f"✅ Pitch deck fully rewritten with {provider_name}")
                    except Exception as e:
                        error_str = str(e)
                        # Check if it's a quota error - try fallback
                        if "429" in error_str or "quota" in error_str.lower() or "ResourceExhausted" in error_str:
                            logger.warning(f"⚠️ {provider_name} quota exceeded. Trying fallback to OpenAI...")
                            # Try OpenAI as fallback
                            if provider_name == "Gemini" and use_openai:
                                try:
                                    from .openai_enhancer import OpenAISlideEnhancer
                                    openai_enhancer = OpenAISlideEnhancer()
                                    if openai_enhancer.enabled:
                                        logger.info("🔄 Falling back to ChatGPT for rewrite...")
                                        slides = openai_enhancer.rewrite_full_pitch_deck(slides, company_name, company_data)
                                        logger.info("✅ Pitch deck rewritten with ChatGPT (fallback)")
                                    else:
                                        logger.warning("⚠️ OpenAI also not available. Continuing without AI enhancement.")
                                except Exception as fallback_error:
                                    logger.warning(f"⚠️ Fallback to OpenAI also failed: {fallback_error}")
                                    logger.info("ℹ️ Continuing with original slides (no AI enhancement)")
                        else:
                            logger.warning(f"{provider_name} rewrite failed: {e}. Continuing without enhancement.")
                else:
                    logger.info("⚠️ No AI provider available (check API keys)")
            
            # 🆕 INDIVIDUAL SLIDE ENHANCEMENT (backup if full rewrite not used)
            elif enhance_with_ai:
                enhancer = None
                provider_name = None
                
                use_gemini = ai_provider.lower() in ["auto", "gemini"]
                use_openai = ai_provider.lower() in ["auto", "openai"]
                
                if use_gemini:
                    try:
                        from .gemini_enhancer import GeminiSlideEnhancer
                        gemini_enhancer = GeminiSlideEnhancer()
                        if gemini_enhancer.enabled:
                            enhancer = gemini_enhancer
                            provider_name = "Gemini"
                    except Exception:
                        pass
                
                if not enhancer and use_openai:
                    try:
                        from .openai_enhancer import OpenAISlideEnhancer
                        openai_enhancer = OpenAISlideEnhancer()
                        if openai_enhancer.enabled:
                            enhancer = openai_enhancer
                            provider_name = "ChatGPT"
                    except Exception:
                        pass
                
                if enhancer:
                    try:
                        logger.info(f"🤖 Enhancing slides with {provider_name}...")
                        slides = enhancer.enhance_slides_batch(slides, company_name, enhance_all=True)
                        logger.info(f"✅ Slides enhanced with {provider_name}")
                    except Exception as e:
                        logger.warning(f"{provider_name} enhancement failed, continuing without: {e}")
                else:
                    logger.info("⚠️ No AI provider available (check API keys)")
            
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            
            logger.info(f"Exporting {len(slides)} slides to PowerPoint (professional)")
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Premium color scheme - modern and professional
            primary_color = RGBColor(30, 58, 138)  # Deep blue
            secondary_color = RGBColor(15, 23, 42)  # Navy
            accent_color = RGBColor(59, 130, 246)  # Bright blue
            highlight_color = RGBColor(34, 197, 94)  # Green accent
            text_color = RGBColor(30, 41, 59)  # Slate
            light_bg = RGBColor(248, 250, 252)  # Light gray background
            white = RGBColor(255, 255, 255)
            
            # Initialize image generators
            image_fetcher = None
            gemini_image_generator = None
            stitched_pptx_path = None  # Will store stitched PPTX if full slide images are generated
            
            if include_images:
                # Try Gemini FULL SLIDE IMAGE generation first (if using Gemini)
                use_gemini = ai_provider.lower() in ["auto", "gemini"]
                if use_gemini:
                    try:
                        from .gemini_image_generator import GeminiImageGenerator
                        gemini_image_generator = GeminiImageGenerator()
                        if gemini_image_generator.enabled:
                            logger.info("🎨 Using Gemini to generate FULL SLIDE IMAGES (one image per slide)...")
                            # Generate full slide images for all slides
                            gemini_images = gemini_image_generator.generate_images_for_slides(
                                slides, company_name, company_data
                            )
                            
                            # Check if we got actual images (not just descriptions)
                            # Filter out None, BytesIO objects, and ensure paths exist
                            generated_images = {}
                            for k, v in gemini_images.items():
                                if v is None:
                                    continue
                                # Skip BytesIO objects
                                if hasattr(v, 'read'):
                                    logger.warning(f"⚠️ Slide {k} has BytesIO object instead of file path - trying to find saved file...")
                                    # Try to find the saved image file
                                    saved_images = list(self.output_dir.glob(f"slide_{k}_gemini_*.png"))
                                    if saved_images:
                                        # Use the most recent one
                                        image_path = str(max(saved_images, key=lambda p: p.stat().st_mtime))
                                        logger.info(f"✅ Found saved image for slide {k}: {image_path}")
                                        generated_images[k] = image_path
                                    continue
                                # Ensure it's a string path
                                image_path = str(v)
                                # Check if file exists
                                if Path(image_path).exists():
                                    generated_images[k] = image_path
                                else:
                                    logger.warning(f"⚠️ Image file not found for slide {k}: {image_path}")
                                    # Try to find the saved image file
                                    saved_images = list(Path("exports/gemini_images").glob(f"slide_{k}_gemini_*.png"))
                                    if saved_images:
                                        # Use the most recent one
                                        image_path = str(max(saved_images, key=lambda p: p.stat().st_mtime))
                                        logger.info(f"✅ Found saved image for slide {k}: {image_path}")
                                        generated_images[k] = image_path
                            
                            if generated_images and len(generated_images) > 0:
                                logger.info(f"✅ Generated {len(generated_images)} full slide images - stitching into PowerPoint...")
                                
                                # 🆕 STITCH IMAGES INTO PPT
                                try:
                                    from .slide_stitcher import SlideStitcher
                                    stitcher = SlideStitcher()
                                    stitched_pptx = stitcher.stitch_images_to_pptx(generated_images, company_name)
                                    if stitched_pptx:
                                        logger.info(f"✅ Stitched {len(generated_images)} full slide images into PowerPoint: {stitched_pptx}")
                                        stitched_pptx_path = stitched_pptx
                                        # If we have stitched PPTX, return it early (skip normal PPT generation)
                                        return stitched_pptx
                                    else:
                                        logger.warning("⚠️ Could not stitch images into PowerPoint")
                                except Exception as e:
                                    logger.error(f"Error stitching images: {e}", exc_info=True)
                            else:
                                logger.info("⚠️ Gemini image generation created descriptions but not images yet.")
                                logger.info("💡 Images will be generated when Gemini image API is available.")
                                logger.info("💡 For now, using fallback image fetcher...")
                        else:
                            logger.info("⚠️ Gemini image generation not available, using fallback")
                    except Exception as e:
                        logger.warning(f"Gemini image generation not available: {e}")
                
                # Fallback to regular image fetcher (if no full slide images were generated)
                if not stitched_pptx_path:
                    try:
                        from .image_fetcher import ImageFetcher
                        image_fetcher = ImageFetcher()
                    except Exception as e:
                        logger.warning(f"Image fetcher not available: {e}")
                        include_images = False
            
            # 🎨 ENHANCED Premium title slide design with background image
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # 🆕 Try to get a background image for title slide
            title_bg_image = None
            if include_images and image_fetcher:
                try:
                    # Get a professional background image for the company
                    title_bg_image = image_fetcher.get_image_for_slide(
                        f"{company_name} professional background",
                        f"{company_name} modern business",
                        ["professional", "business", "modern", "background"]
                    )
                    if title_bg_image and not Path(title_bg_image).exists():
                        title_bg_image = None
                except Exception as e:
                    logger.debug(f"Could not fetch title background image: {e}")
            
            # 🆕 Enhanced background with image or gradient
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height
            
            if title_bg_image and Path(title_bg_image).exists():
                # Use background image with overlay
                try:
                    bg_img = slide.shapes.add_picture(title_bg_image, left, top, width, height)
                    # Add dark overlay for text readability
                    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark overlay
                    overlay.fill.transparency = 0.6  # 60% transparency
                    overlay.line.fill.background()
                    # Send overlay to back (but above image)
                    slide.shapes._spTree.remove(overlay._element)
                    slide.shapes._spTree.insert(2, overlay._element)
                    logger.info("Added background image to title slide")
                except Exception as e:
                    logger.warning(f"Could not add background image: {e}")
                    # Fallback to solid color
                    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                    background.fill.solid()
                    background.fill.fore_color.rgb = secondary_color
                    background.line.fill.background()
            else:
                # Enhanced gradient effect (using layered shapes for depth)
                background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                background.fill.solid()
                background.fill.fore_color.rgb = secondary_color
                background.line.fill.background()
                
                # 🆕 Add subtle gradient overlay effect
                gradient_overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                gradient_overlay.fill.solid()
                gradient_overlay.fill.fore_color.rgb = RGBColor(30, 58, 138)  # Lighter blue
                gradient_overlay.fill.transparency = 0.3
                gradient_overlay.line.fill.background()
                slide.shapes._spTree.remove(gradient_overlay._element)
                slide.shapes._spTree.insert(2, gradient_overlay._element)
            
            # 🆕 Enhanced accent bar with gradient effect
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.4))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = accent_color
            accent_bar.line.fill.background()
            
            # 🆕 Add subtle shine effect on accent bar
            shine = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.1))
            shine.fill.solid()
            shine.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shine.fill.transparency = 0.4
            shine.line.fill.background()
            
            # Company name - large and bold
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(1.2)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = company_name
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(56)
            title_para.font.bold = True
            title_para.font.color.rgb = white
            title_para.alignment = PP_ALIGN.CENTER
            
            # Tagline/subtitle
            top = Inches(4)
            height = Inches(0.6)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = "Pitch Deck"
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(203, 213, 225)  # Light gray
            subtitle_para.alignment = PP_ALIGN.CENTER
            
            # Date
            top = Inches(5)
            height = Inches(0.4)
            date_box = slide.shapes.add_textbox(left, top, width, height)
            date_frame = date_box.text_frame
            date_frame.text = datetime.now().strftime('%B %Y')
            date_para = date_frame.paragraphs[0]
            date_para.font.size = Pt(16)
            date_para.font.color.rgb = RGBColor(148, 163, 184)  # Medium gray
            date_para.alignment = PP_ALIGN.CENTER
            
            # Initialize chart generator for financial data
            chart_generator = None
            available_charts = {}
            if company_data:
                try:
                    from .chart_generator import ChartGenerator
                    chart_generator = ChartGenerator()
                    # Generate all relevant charts upfront
                    available_charts = chart_generator.get_charts_for_slides(company_data)
                    logger.info(f"Generated {len([c for c in available_charts.values() if c])} charts for presentation")
                except Exception as e:
                    logger.warning(f"Chart generator not available: {e}")
            
            # Premium content slides with enhanced design
            for slide_idx, slide_data in enumerate(slides):
                slide = prs.slides.add_slide(blank_layout)
                
                slide_title = slide_data.get('title', 'Untitled')
                slide_content = slide_data.get('content', '')
                key_points = slide_data.get('key_points', [])
                
                # 🎨 ENHANCED: Add subtle background gradient/pattern
                # Create a light gradient background for visual appeal
                bg_left = Inches(0)
                bg_top = Inches(0)
                bg_width = prs.slide_width
                bg_height = prs.slide_height
                background_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bg_left, bg_top, bg_width, bg_height)
                background_shape.fill.solid()
                background_shape.fill.fore_color.rgb = light_bg  # Light background
                background_shape.line.fill.background()
                # Send to back
                slide.shapes._spTree.remove(background_shape._element)
                slide.shapes._spTree.insert(2, background_shape._element)
                
                # Try to get image - first check if already stored in slide data
                image_path = slide_data.get('image_path')
                
                # If not stored, try to fetch it
                if not image_path and include_images and image_fetcher:
                    try:
                        keywords = image_fetcher.get_slide_keywords(slide_title, slide_content)
                        image_path = image_fetcher.get_image_for_slide(slide_title, slide_content, keywords)
                    except Exception as e:
                        logger.warning(f"Could not fetch image: {e}")
                
                # Verify image exists
                if image_path and not Path(image_path).exists():
                    image_path = None
                
                # 🎨 ENHANCED: Premium header with gradient effect and better styling
                left = Inches(0)
                top = Inches(0)
                width = prs.slide_width
                height = Inches(1.3)  # Slightly taller for better presence
                header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                header.fill.solid()
                header.fill.fore_color.rgb = primary_color
                header.line.fill.background()
                
                # Accent line with gradient effect (thicker)
                accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.2))
                accent_line.fill.solid()
                accent_line.fill.fore_color.rgb = accent_color
                accent_line.line.fill.background()
                
                # 🆕 Add subtle pattern/shadow effect for depth
                shadow_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height - Inches(0.05), width, Inches(0.05))
                shadow_bar.fill.solid()
                shadow_bar.fill.fore_color.rgb = RGBColor(10, 20, 50)  # Darker shadow
                shadow_bar.line.fill.background()
                
                # Slide number badge (rounded effect with background)
                left_num = Inches(0.3)
                top_num = Inches(0.35)
                width_num = Inches(0.8)
                height_num = Inches(0.5)
                num_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_num, top_num, width_num, height_num)
                num_bg.fill.solid()
                num_bg.fill.fore_color.rgb = white
                num_bg.line.fill.background()
                
                num_box = slide.shapes.add_textbox(left_num, top_num, width_num, height_num)
                num_frame = num_box.text_frame
                num_frame.text = str(slide_data.get('slide_number', '?'))
                num_para = num_frame.paragraphs[0]
                num_para.font.size = Pt(22)
                num_para.font.bold = True
                num_para.font.color.rgb = primary_color
                num_para.alignment = PP_ALIGN.CENTER
                
                # Slide title - larger and bolder
                left_title = Inches(1.3)
                top_title = Inches(0.3)
                width_title = Inches(7.5)
                height_title = Inches(0.6)
                title_box = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
                title_frame = title_box.text_frame
                title_frame.text = slide_title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = white
                
                # Main content area (larger now, no speech/talking points)
                left_content = Inches(0.5)
                top_content = Inches(1.6)
                width_content = Inches(4.8) if image_path else Inches(9)
                height_content = Inches(5.5)  # Increased height since no speech/talking points
                content_box = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
                content_frame = content_box.text_frame
                content_frame.word_wrap = True
                
                # Add main content with better formatting
                if slide_content:
                    p = content_frame.paragraphs[0]
                    p.text = slide_content
                    p.font.size = Pt(18)  # Slightly larger since more space
                    p.font.color.rgb = text_color
                    p.space_after = Pt(14)
                    p.line_spacing = 1.2
                
                # Add key points as bullets with better styling
                if key_points:
                    for i, point in enumerate(key_points):
                        if i == 0 and not slide_content:
                            p = content_frame.paragraphs[0]
                        else:
                            p = content_frame.add_paragraph()
                        p.text = f"• {point}"
                        p.level = 0
                        p.font.size = Pt(16)  # Slightly larger
                        p.font.color.rgb = text_color
                        p.space_after = Pt(10)
                        p.space_before = Pt(6)
                        p.line_spacing = 1.3
                
                # 🆕 CHECK FOR CHARTS - Add charts to relevant slides automatically
                chart_path = None
                slide_title_lower = slide_title.lower()
                
                # Determine which chart to use based on slide content
                if 'financial' in slide_title_lower or 'revenue' in slide_title_lower or 'traction' in slide_title_lower:
                    chart_path = available_charts.get('revenue') or available_charts.get('projection')
                elif 'growth' in slide_title_lower or 'metrics' in slide_title_lower:
                    chart_path = available_charts.get('growth') or available_charts.get('metrics')
                elif 'projection' in slide_title_lower or 'forecast' in slide_title_lower:
                    chart_path = available_charts.get('projection')
                
                # If chart exists and is valid, use it instead of/alongside image
                if chart_path and Path(chart_path).exists():
                    try:
                        # Place chart prominently (larger than regular images)
                        left_chart = Inches(5.0)
                        top_chart = Inches(1.5)
                        width_chart = Inches(4.8)
                        height_chart = Inches(5.5)
                        
                        chart_img = slide.shapes.add_picture(str(chart_path), left_chart, top_chart, width_chart, height_chart)
                        
                        # Add professional border and shadow effect
                        chart_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                                              left_chart - Inches(0.1), 
                                                              top_chart - Inches(0.1), 
                                                              width_chart + Inches(0.2), 
                                                              height_chart + Inches(0.2))
                        chart_border.fill.background()
                        chart_border.line.color.rgb = RGBColor(200, 200, 200)
                        chart_border.line.width = Pt(2)
                        # Send border to back
                        slide.shapes._spTree.remove(chart_border._element)
                        slide.shapes._spTree.insert(2, chart_border._element)
                        
                        # Adjust content width since chart is taking space
                        content_box.width = Inches(4.5)
                        content_box.left = Inches(0.5)
                        
                        logger.info(f"✅ Added chart to slide: {slide_title}")
                        image_path = None  # Don't add regular image if chart is present
                    except Exception as e:
                        logger.warning(f"Could not add chart to slide: {e}")
                
                # 🎨 ENHANCED: Add image with better styling (if no chart)
                if image_path and not chart_path:
                    # Handle BytesIO objects - convert to file path if needed
                    if hasattr(image_path, 'read'):
                        # It's a BytesIO object, need to save it first
                        from io import BytesIO
                        import tempfile
                        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                        temp_file.write(image_path.read())
                        temp_file.close()
                        image_path = temp_file.name
                        logger.info(f"💾 Saved BytesIO image to temporary file: {image_path}")
                    
                    # Ensure it's a string path and exists
                    image_path = str(image_path)
                    
                    if not Path(image_path).exists():
                        logger.warning(f"⚠️ Image file not found: {image_path}")
                        image_path = None
                    
                    if image_path:
                        try:
                            left_img = Inches(5.5)
                            top_img = Inches(1.6)
                            width_img = Inches(4.2)
                            height_img = Inches(5.7)
                            
                            # Image with border
                            img = slide.shapes.add_picture(image_path, left_img, top_img, width_img, height_img)
                            # Add subtle border effect
                            img_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_img - Inches(0.05), 
                                                              top_img - Inches(0.05), width_img + Inches(0.1), 
                                                              height_img + Inches(0.1))
                            img_line.fill.background()
                            img_line.line.color.rgb = RGBColor(200, 200, 200)
                            img_line.line.width = Pt(1)
                        except Exception as e:
                            logger.warning(f"Could not add image to slide: {e}")
            
            # Save
            filename = f"{company_name.replace(' ', '_')}_pitch_deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            filepath = self.output_dir / filename
            
            prs.save(str(filepath))
            logger.info(f"Professional PowerPoint exported to: {filepath}")
            return str(filepath)
            
        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            return None
        except Exception as e:
            logger.error(f"Error creating PowerPoint: {e}")
            return None

