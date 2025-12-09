"""Chat interface for TechScopeAI using Streamlit."""

import streamlit as st
import streamlit.components.v1 as components
import json
import logging
import os
from pathlib import Path
from typing import Dict, Optional, List
from datetime import datetime
import sys

# Fix tokenizers parallelism warning
os.environ["TOKENIZERS_PARALLELISM"] = "false"

# Load environment variables from .env file
from dotenv import load_dotenv
load_dotenv()

# Add src to path
sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from src.agents.pitch_agent import PitchAgent
from src.agents.competitive_agent import CompetitiveAgent
from src.agents.patent_agent import PatentAgent
from src.agents.policy_agent import PolicyAgent
from src.agents.marketing_agent import MarketingAgent
from src.agents.team_agent import TeamAgent
from src.agents.coordinator_agent import CoordinatorAgent
from src.agents.supervisor_agent import SupervisorAgent
from src.rag.embedder import Embedder
from src.rag.vector_store import VectorStore
from src.rag.retriever import Retriever
from src.data.load_company_data import load_test_company_data, format_company_data_for_pitch, load_company_data
from src.utils.exporters import PitchExporter
from src.utils.user_choices import UserChoiceManager
from src.utils.pitch_question_analyzer import PitchQuestionAnalyzer

logger = logging.getLogger(__name__)

# Configure page
st.set_page_config(
    page_title="TechScopeAI - Pitch Agent",
    page_icon="🚀",
    layout="wide"
)


@st.cache_resource
def load_coordinator_agent():
    """Load and cache CoordinatorAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        agent = CoordinatorAgent(retriever)
        return agent, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
def load_supervisor_agent():
    """Load and cache SupervisorAgent with all agents registered."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Create supervisor
        supervisor = SupervisorAgent(retriever)
        
        # Register all agents
        pitch_agent, _ = load_pitch_agent()
        if pitch_agent:
            supervisor.register_agent("pitch", pitch_agent)
        
        competitive_agent, _ = load_competitive_agent()
        if competitive_agent:
            supervisor.register_agent("competitive", competitive_agent)
        
        patent_agent, _ = load_patent_agent()
        if patent_agent:
            supervisor.register_agent("patent", patent_agent)
        
        policy_agent, _ = load_policy_agent()
        if policy_agent:
            supervisor.register_agent("policy", policy_agent)
        
        marketing_agent, _ = load_marketing_agent()
        if marketing_agent:
            supervisor.register_agent("marketing", marketing_agent)
        
        team_agent, _ = load_team_agent()
        if team_agent:
            supervisor.register_agent("team", team_agent)
        
        coordinator_agent, _ = load_coordinator_agent()
        if coordinator_agent:
            supervisor.register_agent("coordinator", coordinator_agent)
        
        return supervisor, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
def load_pitch_agent():
    """Load and cache PitchAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Initialize agent (will use provider from company_data when generating)
        agent = PitchAgent(retriever, ai_provider="auto")  # Default to auto, will use provider from company_data
        return agent, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
@st.cache_resource(ttl=3600)  # Cache for 1 hour, but will update on code changes
def load_competitive_agent():
    """Load and cache CompetitiveAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Initialize agent with "auto" mode for automatic Gemini fallback on quota errors
        agent = CompetitiveAgent(retriever, model="gpt-4-turbo-preview", ai_provider="auto")
        return agent, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
def load_patent_agent():
    """Load and cache PatentAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Initialize agent (has built-in web search fallback)
        agent = PatentAgent(retriever)
        return agent, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
def load_policy_agent():
    """Load and cache PolicyAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Initialize agent (has built-in web search fallback)
        agent = PolicyAgent(retriever)
        return agent, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
def load_marketing_agent():
    """Load and cache MarketingAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Initialize agent (has built-in web search fallback and image generation)
        agent = MarketingAgent(retriever)
        return agent, None
    except Exception as e:
        return None, str(e)


@st.cache_resource
def load_team_agent():
    """Load and cache TeamAgent."""
    try:
        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None, "OPENAI_API_KEY not found in environment. Add it to .env file"
        
        # Initialize RAG components (PostgreSQL for RAG, Weaviate agents separate)
        embedder = Embedder(use_openai=False)
        embedding_model = embedder.get_embedding_model()
        # VectorStore defaults to Cloud SQL (use use_local=True for local PostgreSQL)
        vector_store = VectorStore(embedding_model=embedding_model)
        
        retriever = Retriever(
            vector_store=vector_store,
            embedder=embedder
        )
        
        # Initialize agent (has built-in web search fallback)
        agent = TeamAgent(retriever)
        return agent, None
    except Exception as e:
        return None, str(e)


def load_company_data_file(company_id: str) -> Optional[Dict]:
    """Load company data from file."""
    company_path = Path(f"src/data/user_companies/{company_id}.json")
    if company_path.exists():
        with open(company_path, 'r') as f:
            return json.load(f)
    return None


def save_company_data(company_id: str, data: Dict):
    """Save company data to file."""
    company_path = Path(f"src/data/user_companies/{company_id}.json")
    company_path.parent.mkdir(parents=True, exist_ok=True)
    with open(company_path, 'w') as f:
        json.dump(data, f, indent=2)


def get_patent_questionnaire():
    """Get patent questionnaire questions."""
    from src.agents.patent_agent import PatentAgent
    return PatentAgent.PATENT_QUESTIONNAIRE


def get_marketing_questionnaire():
    """Get marketing content questionnaire questions."""
    from src.agents.marketing_agent import MarketingAgent
    return MarketingAgent.MARKETING_QUESTIONNAIRE


def get_policy_questionnaire():
    """Get policy questionnaire questions."""
    from src.agents.policy_agent import PolicyAgent
    return PolicyAgent.POLICY_QUESTIONNAIRE


def get_team_questionnaire():
    """Get team building questionnaire questions."""
    from src.agents.team_agent import TeamAgent
    return TeamAgent.TEAM_QUESTIONNAIRE


def get_role_specific_questionnaire():
    """Get role-specific questionnaire questions."""
    from src.agents.team_agent import TeamAgent
    return TeamAgent.ROLE_SPECIFIC_QUESTIONNAIRE


def should_show_question(question, answers):
    """Check if a question should be shown based on dependencies."""
    if "depends_on" not in question:
        return True
    
    depends = question["depends_on"]
    for key, value in depends.items():
        if key not in answers:
            return False
        answer = answers[key]
        if isinstance(value, list):
            if answer not in value:
                return False
        else:
            if answer != value:
                return False
    return True


def get_pitch_questions_for_company(company_data: Dict):
    """Get intelligent pitch questions for company based on existing data."""
    analyzer = PitchQuestionAnalyzer()
    questions_to_ask, existing_data = analyzer.analyze_company_data(company_data)
    return questions_to_ask, existing_data


def generate_pitch_options(agent: PitchAgent, company_data: Dict, num_options: int = 3) -> List[Dict]:
    """
    Generate multiple pitch deck options with different angles/styles.
    
    Args:
        agent: PitchAgent instance
        company_data: Company data dictionary
        num_options: Number of pitch options to generate (default: 3)
        
    Returns:
        List of pitch options (each with slides, title, description)
    """
    options = []
    
    # Different pitch angles/styles
    pitch_angles = [
        {"style": "problem_solution", "focus": "Emphasize problem pain and solution fit", "tone": "urgent"},
        {"style": "market_opportunity", "focus": "Highlight market size and growth potential", "tone": "ambitious"},
        {"style": "traction_focused", "focus": "Showcase existing traction and validation", "tone": "confident"}
    ]
    
    # Limit to requested number
    pitch_angles = pitch_angles[:num_options]
    
    for i, angle in enumerate(pitch_angles, 1):
        # Modify company_data to include pitch angle guidance
        enhanced_data = {
            **company_data,
            "pitch_angle": angle["style"],
            "pitch_focus": angle["focus"],
            "pitch_tone": angle["tone"],
            "pitch_option_number": i
        }
        
        # Generate slides with this angle
        try:
            slides_data = agent.generate_slides(enhanced_data, gamma_only=True)  # Don't generate PPT yet
            
            option = {
                "option_number": i,
                "title": f"Option {i}: {angle['focus']}",
                "description": f"{angle['tone'].title()} tone - {angle['focus']}",
                "style": angle["style"],
                "slides_data": slides_data,
                "slides_count": slides_data.get("total_slides", 0)
            }
            options.append(option)
        except Exception as e:
            logger.error(f"Error generating pitch option {i}: {e}")
            # Continue with other options
    
    return options


def show_choice_selection(choice_manager: UserChoiceManager, agent_name: str, 
                          action: str, user_id: str, company_id: str) -> Optional[Dict]:
    """
    Show choice selection UI and return selected choice.
    
    Args:
        choice_manager: UserChoiceManager instance
        agent_name: Name of the agent
        action: Action name
        user_id: User identifier
        company_id: Company identifier
        
    Returns:
        Selected choice dictionary or None if not yet selected
    """
    choices = choice_manager.get_choices_for_action(agent_name, action)
    
    if len(choices) == 1:
        # Only one choice, auto-select it
        return choices[0]
    
    # Check if a choice was already selected in this session
    choice_key_base = f"selected_choice_{agent_name}_{action}"
    if choice_key_base in st.session_state:
        return st.session_state[choice_key_base]
    
    # Show choice selection UI
    st.markdown("### 🎯 Choose an Option")
    st.info("Please select how you'd like to proceed:")
    
    # Display choices in columns
    num_cols = min(len(choices), 3)
    cols = st.columns(num_cols)
    
    for idx, choice in enumerate(choices):
        with cols[idx % num_cols]:
            with st.container():
                st.markdown(f"#### {choice.get('icon', '📌')} {choice.get('label', 'Option')}")
                st.caption(choice.get('description', ''))
                
                choice_key = f"choice_btn_{agent_name}_{action}_{choice['id']}"
                if st.button(f"Select", key=choice_key, use_container_width=True):
                    # Store selected choice in session state
                    st.session_state[choice_key_base] = choice
                    st.rerun()
    
    return None


def execute_action_with_choice(choice_manager: UserChoiceManager, agent_name: str,
                                action: str, choice: Dict, user_id: str, company_id: str,
                                execute_func, *args, **kwargs):
    """
    Execute an action with the selected choice and save the choice.
    
    Args:
        choice_manager: UserChoiceManager instance
        agent_name: Name of the agent
        action: Action name
        choice: Selected choice dictionary
        user_id: User identifier
        company_id: Company identifier
        execute_func: Function to execute the action
        *args, **kwargs: Arguments to pass to execute_func
        
    Returns:
        Result from execute_func
    """
    # Save the choice
    choice_manager.save_choice(
        user_id=user_id,
        company_id=company_id,
        agent_name=agent_name,
        action=action,
        choice_id=choice['id'],
        choice_data=choice,
        context={
            "timestamp": datetime.now().isoformat(),
            "action_args": str(args),
            "action_kwargs": {k: str(v) for k, v in kwargs.items()}
        }
    )
    
    # Merge choice value into kwargs
    if 'choice_value' not in kwargs:
        kwargs['choice_value'] = choice.get('value', {})
    else:
        kwargs['choice_value'].update(choice.get('value', {}))
    
    # Execute the function
    return execute_func(*args, **kwargs)


def main():
    """Main chat interface."""
    st.title("🚀 TechScopeAI - Pitch Agent")
    st.markdown("**Your AI-powered pitch deck advisor**")
    
    # Initialize session state
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'company_data' not in st.session_state:
        st.session_state.company_data = None
    if 'company_id' not in st.session_state:
        st.session_state.company_id = None
    if 'slides' not in st.session_state:
        st.session_state.slides = None
    if 'current_slide' not in st.session_state:
        st.session_state.current_slide = 0
    if 'elevator_pitch' not in st.session_state:
        st.session_state.elevator_pitch = None
    if 'use_supervisor' not in st.session_state:
        st.session_state.use_supervisor = True  # Enable supervisor by default
    if 'pending_action' not in st.session_state:
        st.session_state.pending_action = None  # Store pending action waiting for choice
    if 'pending_action_data' not in st.session_state:
        st.session_state.pending_action_data = None  # Store action context
    if 'user_id' not in st.session_state:
        import uuid
        st.session_state.user_id = str(uuid.uuid4())  # Generate unique user ID for session
    
    # Initialize UserChoiceManager
    choice_manager = UserChoiceManager()
    
    # Load agent
    agent, error = load_pitch_agent()
    
    if error:
        st.error(f"❌ Error loading agent: {error}")
        st.info("💡 Make sure you've:")
        st.info("1. Processed the pitch data (run: python scripts/processing/process_pitch_data.py)")
        st.info("2. Built the RAG index (run: python scripts/processing/build_rag_index.py)")
        st.info("3. Set OPENAI_API_KEY in .env file")
        return
    
    # Sidebar for company details
    with st.sidebar:
        st.header("📋 Company Details")
        
        # Database Connection Status
        st.markdown("---")
        st.subheader("🔌 Database Status")
        try:
            from src.rag.db_config import get_database_url
            from src.rag.vector_store import VectorStore
            from src.rag.embeddings import EmbeddingModel
            import os
            
            # Check if Cloud SQL password is set (indicates Cloud SQL is default)
            cloud_password = os.getenv("CLOUD_SQL_PASSWORD")
            cloud_host = os.getenv("CLOUD_SQL_HOST", "localhost")
            cloud_port = os.getenv("CLOUD_SQL_PORT", "5433")
            
            # Get database URL to check connection type
            db_url = get_database_url()
            
            # Detect connection type
            is_cloud = (
                cloud_password is not None or  # Cloud SQL password set
                ":5433" in db_url or  # Cloud SQL proxy port
                "/cloudsql/" in db_url or  # Cloud SQL socket path
                "cloudsql" in db_url.lower()
            )
            is_local = ":5432" in db_url and not is_cloud
            
            if is_cloud:
                st.success("✅ **Cloud SQL** (GCP)")
                conn_info = db_url.split('@')[-1] if '@' in db_url else f"{cloud_host}:{cloud_port}"
                st.caption(f"📍 {conn_info}")
                
                # Test connection and show collection counts
                try:
                    embedding_model = EmbeddingModel()
                    test_vs = VectorStore(embedding_model=embedding_model)
                    conn = test_vs._get_connection()
                    test_vs._put_connection(conn)
                    
                    # Get collection counts
                    collections = ['competitors_corpus', 'marketing_corpus', 'ip_policy_corpus', 
                                 'policy_corpus', 'job_roles_corpus', 'pitch_examples_corpus']
                    total_docs = 0
                    collection_counts = {}
                    for coll in collections:
                        try:
                            count = test_vs.get_collection_count(coll)
                            total_docs += count
                            if count > 0:
                                collection_counts[coll.replace('_corpus', '')] = count
                        except:
                            pass
                    
                    if total_docs > 0:
                        st.caption(f"📊 **{total_docs:,}** documents indexed")
                        if collection_counts:
                            coll_list = ", ".join([f"{k}: {v:,}" for k, v in list(collection_counts.items())[:3]])
                            st.caption(f"   ({coll_list})")
                    else:
                        st.caption("⚠️ No documents indexed yet")
                    st.caption("✅ Connection active")
                except Exception as e:
                    st.warning(f"⚠️ Connection test failed")
                    st.caption(f"   {str(e)[:60]}...")
            elif is_local:
                st.warning("⚠️ **Local PostgreSQL**")
                conn_info = db_url.split('@')[-1] if '@' in db_url else "localhost:5432"
                st.caption(f"📍 {conn_info}")
                st.caption("💡 Set CLOUD_SQL_PASSWORD in .env to use Cloud SQL")
            else:
                st.info(f"🔗 **Custom Connection**")
                conn_info = db_url.split('@')[-1] if '@' in db_url else "Custom"
                st.caption(f"📍 {conn_info}")
        except Exception as e:
            st.error(f"❌ Status check failed")
            st.caption(f"   {str(e)[:60]}...")
        
        # Supervisor toggle (for debugging)
        st.markdown("---")
        st.subheader("⚙️ Settings")
        use_supervisor_setting = st.checkbox(
            "Use Supervisor Agent (Smart Routing)", 
            value=st.session_state.get("use_supervisor", True),
            help="Enable intelligent query routing via Supervisor Agent. Disable to use direct routing."
        )
        st.session_state.use_supervisor = use_supervisor_setting
        
        # Gamma and Canva removed - not working
        
        # Company ID input
        company_id = st.text_input("Company ID", value=st.session_state.company_id or "default")
        st.session_state.company_id = company_id
        
        # Load test data option
        if st.button("📋 Load Test Company Data"):
            test_data = load_test_company_data()
            if test_data:
                formatted = format_company_data_for_pitch(test_data)
                st.session_state.company_data = formatted
                st.success(f"✅ Loaded test data for: {formatted.get('company_name', 'Unknown')}")
            else:
                st.error("Test company data not found!")
        
        # Load existing company data
        if st.button("📂 Load Saved Company Data"):
            data = load_company_data_file(company_id)
            if data:
                formatted = format_company_data_for_pitch(data)
                st.session_state.company_data = formatted
                st.success("✅ Company data loaded!")
            else:
                st.warning("No company data found. Fill in the form below or load test data.")
        
        # Company details form
        st.subheader("Enter Company Details")
        
        company_name = st.text_input("Company Name", value=st.session_state.company_data.get('company_name', '') if st.session_state.company_data else '')
        industry = st.text_input("Industry", value=st.session_state.company_data.get('industry', '') if st.session_state.company_data else '')
        problem = st.text_area("Problem", value=st.session_state.company_data.get('problem', '') if st.session_state.company_data else '', height=100)
        solution = st.text_area("Solution", value=st.session_state.company_data.get('solution', '') if st.session_state.company_data else '', height=100)
        market = st.text_input("Target Market", value=st.session_state.company_data.get('target_market', '') if st.session_state.company_data else '')
        stage = st.selectbox("Current Stage", ["Pre-Seed", "Seed", "Series A", "Series B", "Series C+"], 
                           index=0 if not st.session_state.company_data else 
                           ["Pre-Seed", "Seed", "Series A", "Series B", "Series C+"].index(st.session_state.company_data.get('current_stage', 'Pre-Seed')) if st.session_state.company_data.get('current_stage') in ["Pre-Seed", "Seed", "Series A", "Series B", "Series C+"] else 0)
        traction = st.text_area("Traction", value=st.session_state.company_data.get('traction', '') if st.session_state.company_data else '', height=80)
        funding_goal = st.text_input("Funding Goal", value=st.session_state.company_data.get('funding_goal', '') if st.session_state.company_data else '')
        
        if st.button("💾 Save Company Data"):
            company_data = {
                'company_name': company_name,
                'industry': industry,
                'problem': problem,
                'solution': solution,
                'target_market': market,
                'current_stage': stage,
                'traction': traction,
                'funding_goal': funding_goal
            }
            save_company_data(company_id, company_data)
            st.session_state.company_data = company_data
            st.success("✅ Company data saved!")
    
    # Main chat area
    st.header("💬 Chat with Pitch Agent")
    
    # Display chat messages
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            
            # Display image if available
            if "image_path" in message and message["image_path"]:
                image_path = message["image_path"]
                if Path(image_path).exists():
                    try:
                        from PIL import Image
                        img = Image.open(image_path)
                        caption = "Generated Instagram Image - Ready to Post!" if message.get('image_generated') else "Generated Marketing Image"
                        st.image(img, caption=caption, width='stretch')
                    except Exception as e:
                        st.warning(f"Could not display image: {e}")
            
            if "sources" in message and message["sources"]:
                # Determine if sources include web search results
                has_web_search = any(source.get('is_web_search', False) for source in message["sources"])
                source_label = "📚 Sources" + (" (Cloud SQL + Web Search)" if has_web_search else " (from Cloud SQL)")
                
                with st.expander(source_label):
                    for i, source in enumerate(message["sources"][:5], 1):  # Show top 5 sources
                        source_name = source.get('source', 'Unknown')
                        source_title = source.get('title', '')
                        snippet = source.get('snippet', '')
                        similarity = source.get('similarity', 0)
                        metadata = source.get('metadata', {})
                        is_web_search = source.get('is_web_search', False)
                        
                        # Display title or source name
                        if source_title:
                            display_title = source_title
                        else:
                            display_title = source_name
                        
                        # Show collection name if available in metadata (for RAG sources)
                        collection_hint = ""
                        if isinstance(metadata, dict) and not is_web_search:
                            if any(coll in source_name.lower() for coll in ['corpus', 'competitor', 'marketing', 'pitch', 'policy', 'job', 'ip']):
                                collection_hint = f" [{source_name}]"
                            elif 'agent' in metadata:
                                collection_hint = f" [{metadata.get('agent', '')}_corpus]"
                        
                        # Display source with clickable link if it's a URL
                        if is_web_search and source_name.startswith('http'):
                            st.markdown(f"**{i}.** [{display_title}]({source_name})")
                        else:
                            st.markdown(f"**{i}.** {display_title}{collection_hint}")
                        
                        # Show snippet if available (for web search results)
                        if snippet and is_web_search:
                            st.markdown(f"   *{snippet[:200]}{'...' if len(snippet) > 200 else ''}*")
                            st.markdown(f"   🔗 [View full article]({source_name})")
                        elif source_name.startswith('http') and not is_web_search:
                            st.markdown(f"   🔗 [View source]({source_name})")
                        
                        # Show similarity and source type
                        if similarity > 0:
                            source_type = "Web Search" if is_web_search else "Cloud SQL"
                            st.caption(f"   Similarity: {similarity:.3f} | {source_type}")
                        else:
                            source_type = "Web Search" if is_web_search else "Cloud SQL Collection"
                            st.caption(f"   {source_type}")
                        
                        if metadata.get('source') and not is_web_search:
                            st.caption(f"   Original source: {metadata.get('source', 'N/A')}")
    
    # Chat input
    if prompt := st.chat_input("Ask about pitch decks, patents, competitors, or generate a pitch..."):
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Generate response
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                try:
                    # Use Supervisor Agent for routing (if enabled)
                    use_supervisor = st.session_state.get("use_supervisor", True)
                    
                    # Try supervisor first if enabled
                    supervisor_success = False
                    if use_supervisor:
                        # Load supervisor and coordinator
                        supervisor, supervisor_error = load_supervisor_agent()
                        coordinator, coordinator_error = load_coordinator_agent()
                        
                        if supervisor_error or coordinator_error:
                            logger.warning(f"Supervisor/Coordinator not available: {supervisor_error or coordinator_error}")
                            use_supervisor = False
                        else:
                            try:
                                # Get company ID
                                company_id = st.session_state.get("company_id", "default")
                                
                                # Store company data in coordinator
                                if st.session_state.company_data:
                                    coordinator.store_company_data(company_id, st.session_state.company_data)
                                
                                # Get context from coordinator
                                context_from_coordinator = coordinator.get_context_for_agent(
                                    company_id, 
                                    "supervisor",
                                    prompt
                                )
                                
                                # Add coordinator context to query context
                                query_context = {
                                    **(st.session_state.company_data or {}),
                                    "coordinator_context": context_from_coordinator.get("summary", ""),
                                    "company_id": company_id
                                }
                                
                                # Route and execute via supervisor
                                response = supervisor.process_query(prompt, query_context)
                                
                                # Check for errors
                                if response.get("error"):
                                    logger.warning(f"Supervisor agent error: {response.get('error')}")
                                    use_supervisor = False
                                elif not response.get("response") or response.get("response") == "No response":
                                    logger.warning("No response from supervisor agent")
                                    use_supervisor = False
                                else:
                                    # Success - store and display
                                    routed_to = response.get("routed_to", "unknown")
                                    coordinator.store_generated_content(
                                        company_id,
                                        routed_to,
                                        "response",
                                        response
                                    )
                                    
                                    # Display response
                                    st.markdown(response.get("response", "No response"))
                                    
                                    # Show routing info
                                    if response.get("routed_to"):
                                        st.caption(f"🤖 Routed to: {response['routed_to']} agent (confidence: {response.get('routing_confidence', 0):.2f})")
                                    
                                    # Show sources
                                    sources = response.get("sources", [])
                                    if sources:
                                        has_web_search = any(s.get('is_web_search', False) for s in sources)
                                        source_label = "📚 Sources" + (" (Cloud SQL + Web Search)" if has_web_search else " (from Cloud SQL)")
                                        
                                        with st.expander(source_label):
                                            for i, source in enumerate(sources[:5], 1):
                                                source_name = source.get('source', 'Unknown')
                                                source_title = source.get('title', '')
                                                snippet = source.get('snippet', '')
                                                similarity = source.get('similarity', 0)
                                                metadata = source.get('metadata', {})
                                                is_web_search = source.get('is_web_search', False)
                                                
                                                display_title = source_title if source_title else source_name
                                                
                                                if is_web_search and source_name.startswith('http'):
                                                    st.markdown(f"**{i}.** [{display_title}]({source_name})")
                                                else:
                                                    st.markdown(f"**{i}.** {display_title}")
                                                
                                                if snippet and is_web_search:
                                                    st.markdown(f"   *{snippet[:200]}{'...' if len(snippet) > 200 else ''}*")
                                                    st.markdown(f"   🔗 [View full article]({source_name})")
                                                
                                                if similarity > 0:
                                                    source_type = "Web Search" if is_web_search else "Cloud SQL"
                                                    st.caption(f"   Similarity: {similarity:.3f} | {source_type}")
                                                else:
                                                    source_type = "Web Search" if is_web_search else "Cloud SQL Collection"
                                                    st.caption(f"   {source_type}")
                                    
                                    st.session_state.messages.append({
                                        "role": "assistant",
                                        "content": response.get("response", ""),
                                        "sources": sources
                                    })
                                    supervisor_success = True
                                    st.rerun()  # Exit here if supervisor succeeded
                            except Exception as e:
                                logger.error(f"Supervisor execution error: {e}", exc_info=True)
                                # Display error to user
                                error_msg = f"⚠️ Supervisor agent encountered an error. Falling back to direct routing...\n\nError: {str(e)}"
                                st.warning(error_msg)
                                use_supervisor = False
                    
                    # Fallback to direct routing if supervisor didn't succeed
                    if not supervisor_success:
                        # Fallback to direct routing (original logic)
                        # Determine mode based on prompt
                        prompt_lower = prompt.lower()
                        
                        # Check if it's a team-related query
                        team_keywords = ["team", "hiring", "job", "recruit", "role", "position", "employee",
                                        "job description", "jd", "salary", "skillset", "candidate", "hire"]
                        is_team_query = any(keyword in prompt_lower for keyword in team_keywords)
                        
                        # Check if it's a marketing-related query
                        marketing_keywords = ["marketing", "instagram", "linkedin", "social media", "content", 
                                             "post", "hashtag", "campaign", "ad", "promotion", "viral",
                                             "engagement", "follower", "trend", "marketing strategy"]
                        is_marketing_query = any(keyword in prompt_lower for keyword in marketing_keywords)
                        
                        # Check if it's a patent-related query
                        patent_keywords = ["patent", "patentability", "prior art", "ip", "intellectual property", 
                                          "filing", "uspto", "patent search", "patent strategy"]
                        is_patent_query = any(keyword in prompt_lower for keyword in patent_keywords)
                        
                        # Check if it's a policy-related query
                        policy_keywords = ["policy", "privacy policy", "terms of service", "tos", "compliance", 
                                          "gdpr", "ccpa", "hipaa", "legal", "regulations", "employee handbook",
                                          "cookie policy", "data protection", "refund policy"]
                        is_policy_query = any(keyword in prompt_lower for keyword in policy_keywords)
                        
                        if is_team_query:
                            # Route to Team Agent
                            team_agent, team_error = load_team_agent()
                            if team_error:
                                st.error(f"Team Agent Error: {team_error}")
                                response = {"response": f"Error: {team_error}", "sources": []}
                            else:
                                if "job description" in prompt_lower or "jd" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    # Would need role details - suggest using questionnaire
                                    response = {"response": "Please use the '👥 Team & Hiring' button to generate a comprehensive job description with all details.", "sources": []}
                                elif "recommend" in prompt_lower or "role" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    team_context = {}  # Would need team info
                                    response = team_agent.analyze_team_needs(context, team_context)
                                else:
                                    response = team_agent.process_query(prompt, st.session_state.company_data)
                        
                        elif is_marketing_query:
                            # Route to Marketing Agent
                            marketing_agent, marketing_error = load_marketing_agent()
                            if marketing_error:
                                st.error(f"Marketing Agent Error: {marketing_error}")
                                response = {"response": f"Error: {marketing_error}", "sources": []}
                            else:
                                if "instagram" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = marketing_agent.generate_instagram_content(context)
                                elif "linkedin" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = marketing_agent.generate_linkedin_content(context)
                                elif "image" in prompt_lower or "generate image" in prompt_lower:
                                    product_desc = prompt if not st.session_state.company_data else st.session_state.company_data.get('solution', prompt)
                                    response = marketing_agent.generate_marketing_image(product_desc, "Professional")
                                elif "strategy" in prompt_lower or "strategies" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = marketing_agent.suggest_marketing_strategies(context)
                                else:
                                    response = marketing_agent.process_query(prompt, st.session_state.company_data)
                        
                        elif is_policy_query:
                            # Route to Policy Agent
                            policy_agent, policy_error = load_policy_agent()
                            if policy_error:
                                st.error(f"Policy Agent Error: {policy_error}")
                                response = {"response": f"Error: {policy_error}", "sources": []}
                            else:
                                if "privacy policy" in prompt_lower or "privacy" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = policy_agent.generate_privacy_policy(context)
                                elif "terms of service" in prompt_lower or "tos" in prompt_lower or "terms" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = policy_agent.generate_terms_of_service(context)
                                elif "compliance" in prompt_lower or "gdpr" in prompt_lower or "ccpa" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = policy_agent.check_compliance(context)
                                elif "hr" in prompt_lower or "employee" in prompt_lower or "handbook" in prompt_lower:
                                    context = st.session_state.company_data or {}
                                    response = policy_agent.generate_hr_policies(context)
                                else:
                                    response = policy_agent.process_query(prompt, st.session_state.company_data)
                        
                        elif is_patent_query:
                            # Route to Patent Agent
                            patent_agent, patent_error = load_patent_agent()
                            if patent_error:
                                st.error(f"Patent Agent Error: {patent_error}")
                                response = {"response": f"Error: {patent_error}", "sources": []}
                            else:
                                # HITL: Check if we need to trigger questionnaire for patent analysis
                                # If user asks about patent but hasn't provided invention details, trigger questionnaire
                                needs_questionnaire = (
                                    ("patent" in prompt_lower or "invention" in prompt_lower or "patentability" in prompt_lower) and
                                    not st.session_state.get("patent_answers") and
                                    ("analyze" in prompt_lower or "help" in prompt_lower or "assess" in prompt_lower or "?" in prompt)
                                )
                                
                                if needs_questionnaire:
                                    # Trigger patent questionnaire automatically (HITL)
                                    if "patent_questionnaire_active" not in st.session_state:
                                        st.session_state.patent_questionnaire_active = False
                                    if "patent_answers" not in st.session_state:
                                        st.session_state.patent_answers = {}
                                    if "patent_question_index" not in st.session_state:
                                        st.session_state.patent_question_index = 0
                                    
                                    st.session_state.patent_questionnaire_active = True
                                    st.session_state.patent_question_index = 0
                                    st.session_state.patent_answers = {}
                                    
                                    # Show message explaining why questionnaire is needed
                                    st.info("🔬 **I need more information to help you with patent analysis.** Please answer a few questions about your invention.")
                                    st.rerun()
                                elif "patentability" in prompt_lower or "assess" in prompt_lower:
                                    invention = prompt if not st.session_state.company_data else f"{st.session_state.company_data.get('solution', prompt)}"
                                    response = patent_agent.assess_patentability(invention, st.session_state.company_data)
                                elif "prior art" in prompt_lower or "search" in prompt_lower:
                                    query = prompt.replace("patent", "").replace("search", "").strip()
                                    response = patent_agent.search_patents(query, st.session_state.company_data)
                                elif "filing" in prompt_lower or "strategy" in prompt_lower:
                                    if st.session_state.company_data:
                                        response = patent_agent.filing_strategy(st.session_state.company_data)
                                    else:
                                        response = patent_agent.process_query(prompt, st.session_state.company_data)
                                else:
                                    response = patent_agent.process_query(prompt, st.session_state.company_data)
                        
                        elif "evaluate" in prompt_lower or "review" in prompt_lower or "improve" in prompt_lower:
                            # Extract pitch text (assume it's after "evaluate:" or similar)
                            pitch_text = prompt
                            if ":" in prompt:
                                pitch_text = prompt.split(":", 1)[1].strip()
                            
                            response = agent.evaluate_pitch(
                                pitch_text=pitch_text,
                                company_context=st.session_state.company_data
                            )
                        elif "outline" in prompt_lower or "generate from outline" in prompt_lower:
                            # Extract outline (simplified - in real app, use structured input)
                            outline = {
                                "sections": ["Problem", "Solution", "Market", "Traction", "Team", "Ask"],
                                "notes": prompt
                            }
                            response = agent.generate_from_outline(
                                outline=outline,
                                company_context=st.session_state.company_data
                            )
                        elif "generate" in prompt_lower and st.session_state.company_data:
                            # Generate from company details
                            response = agent.generate_from_details(st.session_state.company_data)
                        else:
                            # General query - try pitch agent first
                            response = agent.process_query(
                                query=prompt,
                                context=st.session_state.company_data
                            )
                        
                        # Check if this is an image generation response
                        if response.get('success') is not None and 'image_path' in response:
                            # Handle image generation result
                            if response.get('success'):
                                st.success("✅ Image generated successfully!")
                                
                                # Display image
                                image_path = response.get('image_path')
                                if image_path and Path(image_path).exists():
                                    try:
                                        from PIL import Image
                                        img = Image.open(image_path)
                                        st.image(img, caption="Generated Marketing Image", width='stretch')
                                        
                                        # Download button
                                        with open(image_path, 'rb') as f:
                                            st.download_button(
                                                "⬇ Download Image",
                                                f,
                                                file_name=Path(image_path).name,
                                                mime="image/png"
                                            )
                                    except Exception as e:
                                        st.error(f"Could not display image: {e}")
                                
                                # Show image URL if available
                                if response.get('image_url'):
                                    st.markdown(f"**Image URL:** {response.get('image_url')}")
                                
                                # Show prompt used
                                if response.get('prompt_used'):
                                    with st.expander("📝 Prompt Used"):
                                        st.code(response.get('prompt_used'))
                                
                                # Add to messages
                                msg_content = f"✅ Generated marketing image successfully!\n\n"
                                if response.get('image_url'):
                                    msg_content += f"Image URL: {response.get('image_url')}\n"
                                if response.get('image_path'):
                                    msg_content += f"Saved to: {response.get('image_path')}"
                                
                                st.session_state.messages.append({
                                    "role": "assistant",
                                    "content": msg_content,
                                    "image_path": response.get('image_path'),
                                    "image_url": response.get('image_url')
                                })
                            else:
                                # Image generation failed
                                error_msg = f"❌ Image generation failed: {response.get('error', 'Unknown error')}"
                                if response.get('message'):
                                    error_msg += f"\n{response.get('message')}"
                                st.error(error_msg)
                                st.session_state.messages.append({
                                    "role": "assistant",
                                    "content": error_msg
                                })
                        else:
                            # Regular text response
                            response_text = response.get('response', 'No response generated')
                            
                            # Always display response, even if empty
                            if response_text and response_text.strip():
                                st.markdown(response_text)
                            else:
                                st.warning("⚠️ No response generated. Please try rephrasing your question.")
                                response_text = "I couldn't generate a response. Please try asking your question differently or provide more context."
                            
                            # Check if Instagram content has generated image
                            if response.get('image_generated'):
                                if response.get('image_path'):
                                    st.success("✅ Instagram image generated!")
                                    
                                    # Display image
                                    image_path = response.get('image_path')
                                    if image_path and Path(image_path).exists():
                                        try:
                                            from PIL import Image
                                            img = Image.open(image_path)
                                            st.image(img, caption="Generated Instagram Image - Ready to Post!", width='stretch')
                                            
                                            # Download button
                                            with open(image_path, 'rb') as f:
                                                st.download_button(
                                                    "⬇ Download Image for Instagram",
                                                    f,
                                                    file_name=f"instagram_{Path(image_path).name}",
                                                    mime="image/png"
                                                )
                                        except Exception as e:
                                            st.error(f"Could not display image: {e}")
                                            st.info(f"Image path: {image_path}")
                                    else:
                                        st.warning(f"Image file not found at: {response.get('image_path')}")
                                    
                                    # Show image prompt if available
                                    if response.get('image_prompt'):
                                        with st.expander("📝 Image Prompt Used"):
                                            st.code(response.get('image_prompt'))
                                    
                                    # Show image URL if available
                                    if response.get('image_url'):
                                        st.info(f"Image URL: {response.get('image_url')}")
                                else:
                                    st.warning("⚠️ Image generation was attempted but no image path was returned. Check logs for details.")
                            elif response.get('image_generated') is False:
                                st.info("ℹ️ Image generation was not successful. You can still use the text content and manually create an image.")
                            
                            # Show sources
                            if response.get('sources'):
                                sources = response.get('sources', [])
                                has_web_search = any(s.get('is_web_search', False) for s in sources)
                                source_label = "📚 Sources & Citations" + (" (Cloud SQL + Web Search)" if has_web_search else " (from Cloud SQL)")
                                
                                with st.expander(source_label):
                                    for i, source in enumerate(sources[:5], 1):
                                        source_name = source.get('source', 'Unknown')
                                        source_title = source.get('title', '')
                                        snippet = source.get('snippet', '')
                                        similarity = source.get('similarity', 0)
                                        metadata = source.get('metadata', {})
                                        is_web_search = source.get('is_web_search', False)
                                        
                                        display_title = source_title if source_title else source_name
                                        
                                        if is_web_search and source_name.startswith('http'):
                                            st.markdown(f"**{i}.** [{display_title}]({source_name})")
                                        else:
                                            st.markdown(f"**{i}.** {display_title}")
                                        
                                        if snippet and is_web_search:
                                            st.markdown(f"   *{snippet[:200]}{'...' if len(snippet) > 200 else ''}*")
                                            st.markdown(f"   🔗 [View full article]({source_name})")
                                        
                                        if similarity > 0:
                                            source_type = "Web Search" if is_web_search else "Cloud SQL"
                                            st.caption(f"   Similarity: {similarity:.3f} | {source_type}")
                                        else:
                                            source_type = "Web Search" if is_web_search else "Cloud SQL Collection"
                                            st.caption(f"   {source_type}")
                            
                            # Always add to messages, even if response was empty
                            message_content = {
                                "role": "assistant",
                                "content": response_text,  # Use the processed response_text
                                "sources": response.get('sources', [])
                            }
                            
                            # Include image info if available
                            if response.get('image_generated'):
                                message_content['image_generated'] = True
                                message_content['image_path'] = response.get('image_path')
                                message_content['image_url'] = response.get('image_url')
                            
                            st.session_state.messages.append(message_content)
                            st.rerun()
                    
                except Exception as e:
                    error_msg = f"❌ Error: {str(e)}"
                    st.error(error_msg)
                    logger.error(f"Chat error: {e}", exc_info=True)
                    st.session_state.messages.append({
                        "role": "assistant",
                        "content": error_msg
                    })
                    st.rerun()
    
    # Slides Preview Section
    if st.session_state.slides and len(st.session_state.slides.get('slides', [])) > 0:
        st.markdown("---")
        st.subheader("📊 Pitch Deck Slides")
        
        slides_list = st.session_state.slides['slides']
        total_slides = len(slides_list)
        
        # ============================================
        # 🎯 PRIMARY OUTPUT: POWERPOINT PRESENTATION
        # ============================================
        pptx_path = st.session_state.slides.get('pptx_path')
        
        # Fallback: If no path in session, try to find the most recent PowerPoint file
        if not pptx_path or not Path(pptx_path).exists():
            try:
                exports_dir = Path("exports")
                if exports_dir.exists():
                    pptx_files = sorted(exports_dir.glob("*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
                    if pptx_files:
                        # Use most recent PowerPoint file
                        pptx_path = str(pptx_files[0])
                        st.session_state.slides['pptx_path'] = pptx_path
            except Exception:
                pass  # Silently fail
        
        # Always try to generate PowerPoint if slides exist but no PPTX
        # BUT: Skip if we're in Gamma-only mode (pptx_skipped flag)
        is_gamma_only_mode = st.session_state.slides.get('pptx_skipped', False)
        
        if not pptx_path and slides_list and not is_gamma_only_mode:
            print(f"🔄 Auto-generating PowerPoint (not in Gamma-only mode)")
            with st.spinner("🔄 Generating PowerPoint presentation..."):
                try:
                    from src.utils.exporters import PitchExporter
                    exporter = PitchExporter()
                    # Get AI enhancement preference and company data
                    enhance_with_ai = st.session_state.get('enhance_with_ai', True)
                    ai_provider = st.session_state.get('ai_provider', 'auto')  # 🆕 Get AI provider
                    company_data = st.session_state.get('company_data', {})
                    pptx_path = exporter.export_to_powerpoint(
                        slides_list,
                        st.session_state.slides.get('company_name', 'Company'),
                        include_images=True,
                        enhance_with_ai=enhance_with_ai,
                        company_data=company_data,  # 🆕 Pass company data
                        full_rewrite=True,  # 🆕 Enable full AI rewrite
                        ai_provider=ai_provider  # 🆕 Pass AI provider
                    )
                    if pptx_path:
                        st.session_state.slides['pptx_path'] = pptx_path
                        st.success("✅ PowerPoint generated!")
                except Exception as e:
                    st.warning(f"⚠️ Could not auto-generate PowerPoint: {e}")
                    logger.error(f"PowerPoint generation error: {e}", exc_info=True)
        elif is_gamma_only_mode:
            print(f"⏭️ Skipping auto-PowerPoint generation (Gamma-only mode detected)")
            logger.info("Skipping auto-PowerPoint generation - Gamma-only mode")
        
        # Don't show PowerPoint section if we're in Gamma-only mode
        is_gamma_only_mode = st.session_state.slides.get('pptx_skipped', False)
        
        if pptx_path and Path(pptx_path).exists() and not is_gamma_only_mode:
            st.markdown("---")
            st.markdown("## 📊 **YOUR POWERPOINT PRESENTATION**")
            st.markdown("### ✨ AI-Enhanced Professional Pitch Deck")
            
            # Main PowerPoint section
            col_ppt_main1, col_ppt_main2 = st.columns([3, 1])
            with col_ppt_main1:
                st.success(f"✅ **PowerPoint Ready!**")
                st.markdown(f"""
                - **📁 File:** `{Path(pptx_path).name}`
                - **📊 Slides:** {total_slides} professional slides
                - **🎨 Style:** AI-enhanced with modern design
                - **🖼️ Images:** Included (if available)
                """)
                
                # File size
                file_size = Path(pptx_path).stat().st_size / 1024  # KB
                st.caption(f"📦 File size: {file_size:.1f} KB")
            
            with col_ppt_main2:
                with open(pptx_path, 'rb') as f:
                    st.download_button(
                        "⬇️ **DOWNLOAD PPTX**",
                        f.read(),
                        file_name=Path(pptx_path).name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        help="Download your complete PowerPoint presentation",
                        use_container_width=True,
                        type="primary"
                    )
            
            # Preview PowerPoint slides as images
            st.markdown("---")
            st.markdown("### 👁️ **Preview Your Slides**")
            st.info("💡 **This is your main presentation!** Download the PPTX file above to open in PowerPoint, Google Slides, or any presentation software.")
            
            # Try to show slide previews
            try:
                # Convert PPTX to images for preview (if possible)
                from pptx import Presentation
                from PIL import Image
                import io
                
                prs = Presentation(pptx_path)
                num_slides = len(prs.slides)
                
                st.markdown(f"**📑 Your presentation has {num_slides} slides:**")
                
                # Show slide thumbnails/info
                preview_cols = st.columns(min(3, num_slides))
                for i, slide in enumerate(prs.slides[:min(6, num_slides)]):  # Show first 6 slides
                    col_idx = i % 3
                    with preview_cols[col_idx]:
                        # Get slide title or number
                        slide_title = f"Slide {i+1}"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_title = shape.text.strip()[:30] + ("..." if len(shape.text) > 30 else "")
                                break
                        
                        with st.container():
                            st.markdown(f"""
                            <div style="border: 2px solid #3b82f6; border-radius: 8px; padding: 15px; background: linear-gradient(135deg, #eff6ff 0%, #dbeafe 100%); margin-bottom: 10px;">
                                <h4 style="margin: 0; color: #1e40af;">📄 {slide_title}</h4>
                                <p style="margin: 5px 0; color: #64748b; font-size: 0.9em;">Slide {i+1} of {num_slides}</p>
                            </div>
                            """, unsafe_allow_html=True)
                
                if num_slides > 6:
                    st.caption(f"... and {num_slides - 6} more slides. Download the PPTX file to see all slides!")
                    
            except Exception as e:
                # If preview fails, just show info
                st.markdown(f"**✅ Your PowerPoint contains {total_slides} professional slides**")
                st.caption("💡 Download the PPTX file above to view all slides in PowerPoint or Google Slides")
                logger.debug(f"Could not create slide preview: {e}")
            
            st.markdown("---")
        
        # Display Gamma presentation if available
        gamma_result = st.session_state.slides.get('gamma_presentation')
        if gamma_result and (gamma_result.get('success') or gamma_result.get('presentation_url')):
            st.markdown("---")
            st.markdown("## 🎨 **Your Gamma.ai Presentation**")
            gamma_url = (gamma_result.get('presentation_url') or 
                       gamma_result.get('gammaUrl') or 
                       gamma_result.get('url') or
                       gamma_result.get('api_response', {}).get('gammaUrl'))
            
            if gamma_url:
                col_g1, col_g2 = st.columns(2)
                with col_g1:
                    st.markdown("### 🌐 **View Presentation**")
                    st.markdown(f"[**🔗 Open in Gamma.ai →**]({gamma_url})")
                    st.text_input("**Copy URL:**", value=gamma_url, key="gamma_view_slides", disabled=False)
                with col_g2:
                    st.markdown("### ✏️ **Edit Presentation**")
                    edit_url = gamma_result.get('edit_url') or gamma_url
                    st.markdown(f"[**✏️ Edit in Gamma.ai →**]({edit_url})")
                    st.text_input("**Copy URL:**", value=edit_url, key="gamma_edit_slides", disabled=False)
                
                st.info("""
                **📥 How to Download/Export from Gamma:**
                1. Click the **"Open in Gamma.ai"** link above
                2. In Gamma, click the **Share** button (top-right)
                3. Select **Export** from the dropdown
                4. Choose format: **PPTX**, **PDF**, or **PNG**
                5. The file will download automatically
                """)
        
        # Slide navigation
        col_nav1, col_nav2, col_nav3, col_nav4 = st.columns([1, 2, 1, 2])
        with col_nav1:
            if st.button("◀ Previous", disabled=st.session_state.current_slide == 0):
                st.session_state.current_slide = max(0, st.session_state.current_slide - 1)
                st.rerun()
        with col_nav2:
            st.write(f"**Slide {st.session_state.current_slide + 1} of {total_slides}**")
        with col_nav3:
            if st.button("Next ▶", disabled=st.session_state.current_slide >= total_slides - 1):
                st.session_state.current_slide = min(total_slides - 1, st.session_state.current_slide + 1)
                st.rerun()
        with col_nav4:
            # Export buttons
            from src.utils.exporters import PitchExporter
            exporter = PitchExporter()
            col_exp1, col_exp2, col_exp3, col_exp4, col_exp5 = st.columns(5)
            with col_exp1:
                if st.button("📄 Export PDF"):
                    filepath = exporter.export_to_pdf(slides_list, st.session_state.slides.get('company_name', 'Company'))
                    if filepath:
                        with open(filepath, 'rb') as f:
                            st.download_button("⬇ Download PDF", f, file_name=Path(filepath).name, mime="application/pdf")
            with col_exp2:
                if st.button("📊 Export PPTX", help="Generate and download PowerPoint presentation"):
                    with st.spinner("Creating PowerPoint..."):
                        # Get AI enhancement preference (default: True)
                        enhance_with_ai = st.session_state.get('enhance_with_ai', True)
                        ai_provider = st.session_state.get('ai_provider', 'auto')  # 🆕 Get AI provider
                        # Get company data for financial charts and intelligent filtering
                        company_data = st.session_state.get('company_data', {})
                        filepath = exporter.export_to_powerpoint(
                            slides_list, 
                            st.session_state.slides.get('company_name', 'Company'),
                            include_images=True,
                            enhance_with_ai=enhance_with_ai,
                            company_data=company_data,  # 🆕 Pass company data
                            full_rewrite=True,  # 🆕 Enable full AI rewrite
                            ai_provider=ai_provider  # 🆕 Pass AI provider
                        )
                        if filepath:
                            # Update session state with new path
                            st.session_state.slides['pptx_path'] = filepath
                            with open(filepath, 'rb') as f:
                                st.download_button(
                                    "⬇ Download PPTX", 
                                    f.read(), 
                                    file_name=Path(filepath).name, 
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    help="Download the PowerPoint file"
                                )
                            st.success(f"✅ PowerPoint created: {Path(filepath).name}")
                        else:
                            st.error("❌ Failed to create PowerPoint. Check terminal for errors.")
                            st.info("💡 Make sure python-pptx is installed: `pip install python-pptx`")
            with col_exp3:
                if st.button("📝 Export MD"):
                    filepath = exporter.export_to_markdown(slides_list, st.session_state.slides.get('company_name', 'Company'))
                    if filepath:
                        with open(filepath, 'rb') as f:
                            st.download_button("⬇ Download MD", f, file_name=Path(filepath).name, mime="text/markdown")
            # Gamma and Canva removed
        
        # Display current slide
        if 0 <= st.session_state.current_slide < total_slides:
            slide = slides_list[st.session_state.current_slide]
            st.markdown("---")
            with st.container():
                # Create two columns: content on left, image on right
                col_content, col_image = st.columns([2, 1])
                
                with col_content:
                    st.markdown(f"### 📊 Slide {slide.get('slide_number', '?')}: {slide.get('title', 'Untitled')}")
                    
                    # Main content
                    if slide.get('content'):
                        st.markdown("**Content:**")
                        st.markdown(slide.get('content', ''))
                    
                    if slide.get('key_points'):
                        st.markdown("**Key Points:**")
                        for point in slide['key_points']:
                            st.markdown(f"- {point}")
                
                with col_image:
                    # Display image if available
                    image_path = slide.get('image_path')
                    if image_path and Path(image_path).exists():
                        try:
                            from PIL import Image
                            img = Image.open(image_path)
                            st.image(img, caption=slide.get('title', 'Slide Image'), width='stretch')
                        except Exception as e:
                            st.warning(f"Could not display image: {e}")
                    else:
                        st.info("🖼️ No image available for this slide")
                
                # Speech section (in expander - not in PPT)
                if slide.get('speech'):
                    with st.expander("💬 Speech Script (Not in PPT - for reference only)"):
                        st.info(slide.get('speech', ''))
                
                # Talking points section (in expander - not in PPT)
                if slide.get('talking_points'):
                    with st.expander("🎯 Talking Points (Not in PPT - for reference only)"):
                        for point in slide.get('talking_points', []):
                            st.markdown(f"→ **{point}**")
        
        # Slide thumbnails
        st.markdown("---")
        st.subheader("📑 All Slides")
        cols = st.columns(min(5, total_slides))
        for idx, slide in enumerate(slides_list):
            col_idx = idx % 5
            with cols[col_idx]:
                if st.button(f"Slide {idx + 1}\n{slide.get('title', '')[:20]}...", key=f"slide_btn_{idx}"):
                    st.session_state.current_slide = idx
                    st.rerun()
    
    # Elevator Pitch Section (independent of slides)
    if st.session_state.get('elevator_pitch'):
        st.markdown("---")
        st.subheader("🎤 Elevator Pitch")
        
        pitch_data = st.session_state.elevator_pitch
        col_pitch1, col_pitch2 = st.columns([3, 1])
        
        with col_pitch1:
            st.markdown(f"**Duration:** {pitch_data.get('duration_seconds', 60)} seconds")
            st.markdown(f"**Words:** ~{pitch_data.get('estimated_words', 0)} words")
            
            # Display pitch in a nice box
            st.markdown("### Your Elevator Pitch:")
            st.info(pitch_data.get('elevator_pitch', ''))
            
            # Also show in code block for easy copying
            st.markdown("### Copy Text:")
            st.code(pitch_data.get('elevator_pitch', ''), language=None)
        
        with col_pitch2:
            # Copy button info
            st.markdown("### Export")
            pitch_text = pitch_data.get('elevator_pitch', '')
            st.download_button(
                "💾 Download as TXT",
                pitch_text,
                file_name=f"{pitch_data.get('company_name', 'Company')}_ElevatorPitch.txt",
                mime="text/plain"
            )
            st.caption("💡 Tip: Click the code block above and copy to clipboard")
    
    # Show Gamma result persistently (above Quick Actions)
    # Check both gamma_only_result AND slides data
    gamma_result = None
    if st.session_state.get('gamma_only_result'):
        gamma_result = st.session_state.gamma_only_result
    elif st.session_state.get('slides') and st.session_state.slides.get('gamma_presentation'):
        gamma_result = st.session_state.slides.get('gamma_presentation')
    
    if gamma_result and (gamma_result.get('success') or gamma_result.get('presentation_url')):
            st.markdown("---")
            st.markdown("## 🎨 **Your Gamma.ai Presentation**")
            gamma_url = (gamma_result.get('presentation_url') or 
                       gamma_result.get('gammaUrl') or 
                       gamma_result.get('url') or
                       gamma_result.get('api_response', {}).get('gammaUrl'))
            
            if gamma_url:
                col_g1, col_g2 = st.columns(2)
                with col_g1:
                    st.markdown("### 🌐 **View Presentation**")
                    st.markdown(f"[**🔗 Open in Gamma.ai →**]({gamma_url})")
                    st.text_input("**Copy URL:**", value=gamma_url, key="gamma_view_persistent", disabled=False)
                with col_g2:
                    st.markdown("### ✏️ **Edit Presentation**")
                    edit_url = gamma_result.get('edit_url') or gamma_url
                    st.markdown(f"[**✏️ Edit in Gamma.ai →**]({edit_url})")
                    st.text_input("**Copy URL:**", value=edit_url, key="gamma_edit_persistent", disabled=False)
                
                st.info("""
                **📥 How to Download/Export from Gamma:**
                1. Click the **"Open in Gamma.ai"** link above
                2. In Gamma, click the **Share** button (top-right)
                3. Select **Export** from the dropdown
                4. Choose format: **PPTX**, **PDF**, or **PNG**
                5. The file will download automatically
                """)
    
    # Quick actions
    st.markdown("---")
    st.subheader("⚡ Quick Actions")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        if st.button("📝 Generate Pitch"):
            if st.session_state.company_data:
                # Set pending action to show choices
                st.session_state.pending_action = {
                    "agent_name": "pitch",
                    "action": "generate_pitch",
                    "action_type": "generate_pitch"
                }
                st.session_state.pending_action_data = {}
                st.rerun()
            else:
                st.warning("Please enter company details first!")
        
        # Handle pending pitch generation with choice selection
        if st.session_state.pending_action and st.session_state.pending_action.get("action_type") == "generate_pitch":
            company_id = st.session_state.get("company_id", "default")
            user_id = st.session_state.get("user_id", "anonymous")
            
            selected_choice = show_choice_selection(
                choice_manager, 
                "pitch", 
                "generate_pitch",
                user_id,
                company_id
            )
            
            if selected_choice:
                # Execute action with choice
                choice_value = selected_choice.get('value', {})
                format_type = choice_value.get('format', 'text')
                
                with st.spinner("Generating pitch..."):
                    response = agent.generate_from_details(st.session_state.company_data)
                    
                    # Save choice
                    choice_manager.save_choice(
                        user_id=user_id,
                        company_id=company_id,
                        agent_name="pitch",
                        action="generate_pitch",
                        choice_id=selected_choice['id'],
                        choice_data=selected_choice,
                        context={"format_type": format_type}
                    )
                    
                    st.session_state.messages.append({
                        "role": "assistant",
                        "content": response['response'],
                        "sources": response.get('sources', [])
                    })
                    st.session_state.pending_action = None
                    st.session_state.pending_action_data = None
                    st.rerun()
    
    with col2:
        # Quick test button for Gamma/Canva
        if st.button("🧪 Quick Test (3 slides)", help="Test Gamma.ai and Canva with minimal slides"):
            with st.spinner("Creating test presentation with 3 slides..."):
                # Create minimal test slides
                test_slides_data = {
                    "slides": [
                        {
                            "slide_number": 1,
                            "title": "Test Company",
                            "content": "This is a quick test to verify Gamma.ai and Canva integrations work correctly.",
                            "key_points": ["Test integration", "Verify functionality"]
                        },
                        {
                            "slide_number": 2,
                            "title": "Problem",
                            "content": "Testing the integration system.",
                            "key_points": ["Quick test", "Minimal slides"]
                        },
                        {
                            "slide_number": 3,
                            "title": "Solution",
                            "content": "Integration test successful!",
                            "key_points": ["Working", "Ready"]
                        }
                    ],
                    "total_slides": 3,
                    "company_name": "Test Company"
                }
                
                # Gamma and Canva removed
                
                # Generate PowerPoint for test slides too
                try:
                    from src.utils.exporters import PitchExporter
                    exporter = PitchExporter()
                    pptx_path = exporter.export_to_powerpoint(
                        test_slides_data["slides"],
                        "Test Company",
                        include_images=False  # Skip images for quick test
                    )
                    if pptx_path:
                        test_slides_data["pptx_path"] = pptx_path
                        st.success("✅ PowerPoint also generated for test slides!")
                except Exception as e:
                    st.warning(f"⚠️ PowerPoint generation failed: {e}")
                
                # Store test results
                st.session_state.slides = test_slides_data
                st.session_state.current_slide = 0
                st.rerun()
        
        # AI Enhancement Toggle
        col_ai1, col_ai2 = st.columns([2, 1])
        with col_ai1:
            enhance_with_ai = st.checkbox(
                "🤖 Enhance with AI (improves slide content quality)", 
                value=True,
                help="Uses AI to enhance slide titles, content, and key points for better investor appeal"
            )
        with col_ai2:
            # AI Provider Selection (Gemini or OpenAI)
            ai_provider = st.selectbox(
                "AI Provider",
                options=["Auto (Gemini → OpenAI)", "Gemini", "OpenAI"],
                index=0,
                help="Choose AI provider. Auto tries Gemini first, then OpenAI"
            )
        
        # Store in session state so it persists across reruns
        st.session_state['enhance_with_ai'] = enhance_with_ai
        st.session_state['ai_provider'] = ai_provider
        
        # Intelligent Pitch Questionnaire Flow - handled below after patent questionnaire
        
        if st.button("🎯 Generate Slides"):
            if st.session_state.company_data:
                # Check if intelligent questions need to be asked first
                questions_to_ask, existing_data = get_pitch_questions_for_company(st.session_state.company_data)
                
                if questions_to_ask:
                    # Initialize pitch questionnaire if not exists
                    if "pitch_questionnaire_active" not in st.session_state:
                        st.session_state.pitch_questionnaire_active = False
                    if "pitch_answers" not in st.session_state:
                        st.session_state.pitch_answers = {}
                    if "pitch_question_index" not in st.session_state:
                        st.session_state.pitch_question_index = 0
                    if "pitch_existing_data" not in st.session_state:
                        st.session_state.pitch_existing_data = {}
                    
                    # Activate questionnaire
                    st.session_state.pitch_questionnaire_active = True
                    st.session_state.pitch_question_index = 0
                    st.session_state.pitch_answers = {}
                    st.session_state.pitch_existing_data = existing_data
                    st.session_state.pitch_questions_to_ask = questions_to_ask
                    st.info("💡 **Let's build your perfect pitch together!** I'll ask some targeted questions to understand your startup better.")
                else:
                    # All data exists, proceed directly to generation
                    st.session_state.pending_action = {
                        "agent_name": "pitch",
                        "action": "generate_slides",
                        "action_type": "generate_slides",
                        "skip_questionnaire": True  # Flag to skip questionnaire
                    }
                    st.session_state.pending_action_data = {
                        "enhance_with_ai": enhance_with_ai,
                        "ai_provider": ai_provider
                    }
                st.rerun()
            else:
                st.warning("Please enter company details first!")
        
        # Handle pending action with choice selection
        if st.session_state.pending_action:
            pending = st.session_state.pending_action
            if pending.get("action_type") == "generate_slides":
                company_id = st.session_state.get("company_id", "default")
                user_id = st.session_state.get("user_id", "anonymous")
                preferred_choice = pending.get("preferred_choice")
                
                # If there's a preferred choice (e.g., from Gamma button), auto-select it
                if preferred_choice:
                    choices = choice_manager.get_choices_for_action("pitch", "generate_slides")
                    selected_choice = next((c for c in choices if c['id'] == preferred_choice), None)
                    if selected_choice:
                        # Auto-execute with preferred choice
                        choice_value = selected_choice.get('value', {})
                        gamma_mode = choice_value.get('gamma', False)
                        
                        # Show info about AI provider
                        if st.session_state.pending_action_data.get('enhance_with_ai'):
                            ai_prov = st.session_state.pending_action_data.get('ai_provider', 'Auto')
                            if ai_prov == "Gemini":
                                st.info("🤖 Using **Gemini** for AI enhancement. If quota is exceeded, will automatically fallback to OpenAI.")
                            elif ai_prov == "OpenAI":
                                st.info("🤖 Using **ChatGPT (OpenAI)** for AI enhancement.")
                            else:
                                st.info("🤖 Using **Auto mode**: Will try Gemini first, then automatically fallback to OpenAI if needed.")
                        
                        with st.spinner("Generating slides..."):
                            # Include AI provider in company data
                            company_data_with_themes = {
                                **st.session_state.company_data,
                                "enhance_with_ai": st.session_state.pending_action_data.get('enhance_with_ai', True),
                                "ai_provider": st.session_state.pending_action_data.get('ai_provider', 'Auto'),
                                "gamma_theme": st.session_state.get("gamma_theme", "startup-pitch")
                            }
                            
                            # Execute with choice
                            if gamma_mode:
                                slides_data = agent.generate_slides(company_data_with_themes, gamma_only=True)
                            else:
                                slides_data = agent.generate_slides(company_data_with_themes)
                            
                            # Save choice
                            choice_manager.save_choice(
                                user_id=user_id,
                                company_id=company_id,
                                agent_name="pitch",
                                action="generate_slides",
                                choice_id=selected_choice['id'],
                                choice_data=selected_choice,
                                context={"gamma_mode": gamma_mode, "preferred": True}
                            )
                            
                            st.session_state.slides = slides_data
                            st.session_state.current_slide = 0
                            st.session_state.pending_action = None
                            st.session_state.pending_action_data = None
                            st.rerun()
                        return
                
                # Show choice selection UI
                selected_choice = show_choice_selection(
                    choice_manager, 
                    "pitch", 
                    "generate_slides",
                    user_id,
                    company_id
                )
                
                if selected_choice:
                    # Clear the stored choice after using it
                    choice_key_base = f"selected_choice_pitch_generate_slides"
                    if choice_key_base in st.session_state:
                        del st.session_state[choice_key_base]
                    # Execute action with choice
                    choice_value = selected_choice.get('value', {})
                    gamma_mode = choice_value.get('gamma', False)
                    
                    # Show info about AI provider
                    if st.session_state.pending_action_data.get('enhance_with_ai'):
                        ai_prov = st.session_state.pending_action_data.get('ai_provider', 'Auto')
                        if ai_prov == "Gemini":
                            st.info("🤖 Using **Gemini** for AI enhancement. If quota is exceeded, will automatically fallback to OpenAI.")
                        elif ai_prov == "OpenAI":
                            st.info("🤖 Using **ChatGPT (OpenAI)** for AI enhancement.")
                        else:
                            st.info("🤖 Using **Auto mode**: Will try Gemini first, then automatically fallback to OpenAI if needed.")
                    
                    with st.spinner("Generating slides..."):
                        # Include AI provider in company data
                        company_data_with_themes = {
                            **st.session_state.company_data,
                            "enhance_with_ai": st.session_state.pending_action_data.get('enhance_with_ai', True),
                            "ai_provider": st.session_state.pending_action_data.get('ai_provider', 'Auto'),
                            "gamma_theme": st.session_state.get("gamma_theme", "startup-pitch")
                        }
                        
                        # Execute with choice
                        if gamma_mode:
                            slides_data = agent.generate_slides(company_data_with_themes, gamma_only=True)
                        else:
                            slides_data = agent.generate_slides(company_data_with_themes)
                        
                        # Save choice
                        choice_manager.save_choice(
                            user_id=user_id,
                            company_id=company_id,
                            agent_name="pitch",
                            action="generate_slides",
                            choice_id=selected_choice['id'],
                            choice_data=selected_choice,
                            context={"gamma_mode": gamma_mode}
                        )
                        
                        st.session_state.slides = slides_data
                        st.session_state.current_slide = 0
                        st.session_state.pending_action = None
                        st.session_state.pending_action_data = None
                        
                        success_msg = f"✅ Generated {slides_data.get('total_slides', 0)} slides!"
                        
                        # Check what was generated
                        if slides_data.get('pptx_path'):
                            st.success(success_msg)
                        
                        # Show PowerPoint status and download (only if not gamma-only)
                        if slides_data.get('pptx_path') and not gamma_mode:
                            pptx_path = Path(slides_data['pptx_path'])
                            if pptx_path.exists():
                                st.markdown("---")
                                st.success("📊 **PowerPoint Generated Successfully!**")
                                with open(pptx_path, 'rb') as f:
                                    st.download_button(
                                        "⬇️ Download PowerPoint Presentation",
                                        f.read(),
                                        file_name=pptx_path.name,
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        help="Download the complete PowerPoint with all slides",
                                        use_container_width=True
                                    )
                        
                        st.rerun()
        
        # Note: Gamma option is now available through choice selection in "Generate Slides" button
        # This button uses the choice system but auto-selects Gamma for backward compatibility
        if st.button("🚀 Generate Slides Using Gamma", help="Generate slides and create Gamma.ai presentation"):
            if st.session_state.company_data:
                # Set pending action with Gamma as preferred choice (will auto-execute)
                st.session_state.pending_action = {
                    "agent_name": "pitch",
                    "action": "generate_slides",
                    "action_type": "generate_slides",
                    "preferred_choice": "gamma_ppt"  # Auto-select Gamma
                }
                st.session_state.pending_action_data = {
                    "enhance_with_ai": enhance_with_ai,
                    "ai_provider": ai_provider
                }
                st.rerun()
            else:
                st.warning("Please enter company details first!")
        
    with col3:
        if st.button("🎤 Generate Elevator Pitch"):
            if st.session_state.company_data:
                with st.spinner("Crafting your elevator pitch..."):
                    try:
                        # Generate 60-second elevator pitch
                        pitch_result = agent.generate_elevator_pitch(
                            st.session_state.company_data,
                            duration_seconds=60
                        )
                        if pitch_result and pitch_result.get('elevator_pitch'):
                            st.session_state.elevator_pitch = pitch_result
                            st.success("✅ Elevator pitch generated!")
                            st.rerun()
                        else:
                            st.error("❌ Failed to generate elevator pitch. Please try again.")
                            logger.error(f"Elevator pitch generation returned empty result: {pitch_result}")
                    except Exception as e:
                        error_msg = f"Error generating elevator pitch: {str(e)}"
                        st.error(error_msg)
                        logger.error(f"Elevator pitch generation error: {e}", exc_info=True)
            else:
                st.warning("Please enter company details first!")
    
    with col4:
        if st.button("📊 Evaluate Pitch"):
            st.info("Paste your pitch in chat: 'evaluate: [your pitch]'")
    
    # Competitive Agent button
    col_comp = st.columns(1)[0]
    with col_comp:
        if st.button("🔍 Competitive Analysis"):
            if st.session_state.company_data:
                competitive_agent, competitive_error = load_competitive_agent()
                if competitive_error:
                    st.error(f"Error: {competitive_error}")
                else:
                    with st.spinner("Analyzing competitors using your company data..."):
                        try:
                            # Use existing company_data directly - no questionnaire needed
                            competitive_context = st.session_state.company_data.copy()
                            
                            # Perform competitive analysis
                            response = competitive_agent.analyze_competitors(
                                company_details=competitive_context,
                                competitors=None
                            )
                            
                            st.session_state.messages.append({
                                "role": "assistant",
                                "content": response['response'],
                                "sources": response.get('sources', [])
                            })
                            st.rerun()
                        except Exception as e:
                            error_str = str(e)
                            # Check if it's a quota error
                            is_quota_error = ("429" in error_str or "quota" in error_str.lower() or 
                                            "insufficient_quota" in error_str.lower() or 
                                            "RateLimitError" in str(type(e).__name__))
                            
                            if is_quota_error:
                                st.error("""
                                **⚠️ OpenAI Quota Exceeded**
                                
                                Your OpenAI API quota has been exceeded. Options:
                                1. **Wait a few minutes** and try again (if it's rate limiting)
                                2. **Check your OpenAI billing** at https://platform.openai.com/account/billing
                                3. **Add credits** to your OpenAI account
                                4. **Use Gemini instead** - Set `ai_provider` to "Gemini" in settings
                                
                                The system will automatically try to fallback to Gemini if available in "Auto" mode.
                                """)
                                logger.error(f"OpenAI quota error in Competitive Analysis: {e}")
                            else:
                                st.error(f"**Error analyzing competitors:** {error_str}")
                                logger.error(f"Error in Competitive Analysis: {e}", exc_info=True)
            else:
                st.warning("Please enter company details first!")
    
    # Team, Marketing, Policy & Patent Agent buttons
    col_team, col_mkt, col_pol, col_pat = st.columns(4)
    
    with col_team:
        if st.button("👥 Team & Hiring"):
            if "team_questionnaire_active" not in st.session_state:
                st.session_state.team_questionnaire_active = False
            if "team_answers" not in st.session_state:
                st.session_state.team_answers = {}
            if "team_question_index" not in st.session_state:
                st.session_state.team_question_index = 0
            if "team_phase" not in st.session_state:
                st.session_state.team_phase = "analysis"  # "analysis" or "jd_generation"
            if "selected_role" not in st.session_state:
                st.session_state.selected_role = None
            
            st.session_state.team_questionnaire_active = True
            st.session_state.team_question_index = 0
            st.session_state.team_answers = {}
            st.session_state.team_phase = "analysis"
            st.session_state.selected_role = None
            st.rerun()
    
    with col_mkt:
        if st.button("📱 Marketing Content"):
            if "marketing_questionnaire_active" not in st.session_state:
                st.session_state.marketing_questionnaire_active = False
            if "marketing_answers" not in st.session_state:
                st.session_state.marketing_answers = {}
            if "marketing_question_index" not in st.session_state:
                st.session_state.marketing_question_index = 0
            
            st.session_state.marketing_questionnaire_active = True
            st.session_state.marketing_question_index = 0
            st.session_state.marketing_answers = {}
            st.rerun()
    
    with col_pol:
        if st.button("📋 Policy & Compliance"):
            if "policy_questionnaire_active" not in st.session_state:
                st.session_state.policy_questionnaire_active = False
            if "policy_answers" not in st.session_state:
                st.session_state.policy_answers = {}
            if "policy_question_index" not in st.session_state:
                st.session_state.policy_question_index = 0
            
            st.session_state.policy_questionnaire_active = True
            st.session_state.policy_question_index = 0
            st.session_state.policy_answers = {}
            st.rerun()
    
    with col_pat:
        if st.button("🔬 Patent Analysis"):
            if "patent_questionnaire_active" not in st.session_state:
                st.session_state.patent_questionnaire_active = False
            if "patent_answers" not in st.session_state:
                st.session_state.patent_answers = {}
            if "patent_question_index" not in st.session_state:
                st.session_state.patent_question_index = 0
            
            st.session_state.patent_questionnaire_active = True
            st.session_state.patent_question_index = 0
            st.session_state.patent_answers = {}
            st.rerun()
    
    # Marketing Questionnaire Flow
    if "marketing_questionnaire_active" not in st.session_state:
        st.session_state.marketing_questionnaire_active = False
    if "marketing_answers" not in st.session_state:
        st.session_state.marketing_answers = {}
    if "marketing_question_index" not in st.session_state:
        st.session_state.marketing_question_index = 0
    
    if st.session_state.marketing_questionnaire_active:
        st.markdown("---")
        st.subheader("📱 Marketing Content Questionnaire")
        
        questionnaire = get_marketing_questionnaire()
        
        # Filter questions based on dependencies
        visible_questions = [q for q in questionnaire if should_show_question(q, st.session_state.marketing_answers)]
        
        if st.session_state.marketing_question_index < len(visible_questions):
            current_question = visible_questions[st.session_state.marketing_question_index]
            
            st.progress((st.session_state.marketing_question_index + 1) / len(visible_questions))
            st.caption(f"Question {st.session_state.marketing_question_index + 1} of {len(visible_questions)}")
            
            st.markdown(f"### {current_question['question']}")
            
            answer_key = f"marketing_q_{current_question['id']}"
            answer = None
            
            # Render input based on type
            if current_question['type'] == 'textarea':
                answer = st.text_area(
                    "Your answer:",
                    value=st.session_state.marketing_answers.get(current_question['id'], ''),
                    height=150,
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'text':
                answer = st.text_input(
                    "Your answer:",
                    value=st.session_state.marketing_answers.get(current_question['id'], ''),
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'select':
                current_value = st.session_state.marketing_answers.get(current_question['id'])
                default_index = 0
                if current_value and current_value in current_question['options']:
                    default_index = current_question['options'].index(current_value)
                
                answer = st.selectbox(
                    "Select an option:",
                    options=current_question['options'],
                    index=default_index,
                    key=answer_key
                )
            elif current_question['type'] == 'multiselect':
                answer = st.multiselect(
                    "Select all that apply:",
                    options=current_question['options'],
                    default=st.session_state.marketing_answers.get(current_question['id'], []),
                    key=answer_key
                )
            
            # Store answer
            if answer:
                if current_question['type'] == 'multiselect':
                    st.session_state.marketing_answers[current_question['id']] = answer
                else:
                    st.session_state.marketing_answers[current_question['id']] = answer
            
            # Navigation buttons
            col_prev, col_next, col_cancel = st.columns([1, 3, 1])
            
            with col_prev:
                if st.button("◀ Previous", disabled=st.session_state.marketing_question_index == 0, key=f"marketing_prev_{answer_key}"):
                    st.session_state.marketing_question_index -= 1
                    st.rerun()
            
            with col_next:
                is_disabled = False
                if current_question.get('required'):
                    if current_question['type'] == 'multiselect':
                        is_disabled = not answer or len(answer) == 0
                    else:
                        is_disabled = not answer or (isinstance(answer, str) and len(answer.strip()) == 0)
                
                if st.button("Next ▶", disabled=is_disabled, key=f"marketing_next_{answer_key}"):
                    if is_disabled:
                        st.warning("Please answer this question to continue.")
                    else:
                        st.session_state.marketing_question_index += 1
                        st.rerun()
            
            with col_cancel:
                if st.button("❌ Cancel", key=f"marketing_cancel_{answer_key}"):
                    st.session_state.marketing_questionnaire_active = False
                    st.session_state.marketing_answers = {}
                    st.session_state.marketing_question_index = 0
                    st.rerun()
        
        else:
            # All questions answered - generate marketing content
            st.success("✅ All questions answered!")
            
            st.subheader("📋 Your Answers Summary")
            for q in visible_questions:
                if q['id'] in st.session_state.marketing_answers:
                    st.markdown(f"**{q['question']}**")
                    answer = st.session_state.marketing_answers[q['id']]
                    if isinstance(answer, list):
                        st.markdown(f"  → {', '.join(answer)}")
                    else:
                        st.markdown(f"  → {answer}")
                    st.markdown("---")
            
            col_generate, col_restart = st.columns([3, 1])
            
            with col_generate:
                if st.button("📱 Generate Marketing Content", type="primary"):
                    marketing_agent, marketing_error = load_marketing_agent()
                    if marketing_error:
                        st.error(f"Error: {marketing_error}")
                    else:
                        with st.spinner("Generating marketing content..."):
                            # Prepare context from answers
                            marketing_context = st.session_state.marketing_answers.copy()
                            # Merge with company data if available
                            if st.session_state.company_data:
                                marketing_context.update(st.session_state.company_data)
                            
                            # Determine content type and generate
                            platform = marketing_context.get('platform', [])
                            if isinstance(platform, list):
                                platform = platform[0] if platform else 'Both'
                            
                            if 'Instagram' in platform or platform == 'Both':
                                response = marketing_agent.generate_instagram_content(marketing_context)
                            elif 'LinkedIn' in platform:
                                response = marketing_agent.generate_linkedin_content(marketing_context)
                            else:
                                response = marketing_agent.suggest_marketing_strategies(marketing_context)
                            
                            st.session_state.messages.append({
                                "role": "assistant",
                                "content": response['response'],
                                "sources": response.get('sources', []),
                                "image_generated": response.get('image_generated', False),
                                "image_path": response.get('image_path'),
                                "image_url": response.get('image_url')
                            })
                            
                            # Reset questionnaire
                            st.session_state.marketing_questionnaire_active = False
                            st.session_state.marketing_answers = {}
                            st.session_state.marketing_question_index = 0
                            st.rerun()
            
            with col_restart:
                if st.button("🔄 Restart"):
                    st.session_state.marketing_questionnaire_active = False
                    st.session_state.marketing_answers = {}
                    st.session_state.marketing_question_index = 0
                    st.rerun()
    
    # Policy Questionnaire Flow
    if "policy_questionnaire_active" not in st.session_state:
        st.session_state.policy_questionnaire_active = False
    if "policy_answers" not in st.session_state:
        st.session_state.policy_answers = {}
    if "policy_question_index" not in st.session_state:
        st.session_state.policy_question_index = 0
    
    if st.session_state.policy_questionnaire_active:
        st.markdown("---")
        st.subheader("📋 Policy & Compliance Questionnaire")
        
        questionnaire = get_policy_questionnaire()
        
        # Filter questions based on dependencies
        visible_questions = [q for q in questionnaire if should_show_question(q, st.session_state.policy_answers)]
        
        if st.session_state.policy_question_index < len(visible_questions):
            current_question = visible_questions[st.session_state.policy_question_index]
            
            st.progress((st.session_state.policy_question_index + 1) / len(visible_questions))
            st.caption(f"Question {st.session_state.policy_question_index + 1} of {len(visible_questions)}")
            
            st.markdown(f"### {current_question['question']}")
            
            answer_key = f"policy_q_{current_question['id']}"
            answer = None
            
            # Render input based on type
            if current_question['type'] == 'textarea':
                answer = st.text_area(
                    "Your answer:",
                    value=st.session_state.policy_answers.get(current_question['id'], ''),
                    height=150,
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'text':
                answer = st.text_input(
                    "Your answer:",
                    value=st.session_state.policy_answers.get(current_question['id'], ''),
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'select':
                current_value = st.session_state.policy_answers.get(current_question['id'])
                default_index = 0
                if current_value and current_value in current_question['options']:
                    default_index = current_question['options'].index(current_value)
                
                answer = st.selectbox(
                    "Select an option:",
                    options=current_question['options'],
                    index=default_index,
                    key=answer_key
                )
            elif current_question['type'] == 'multiselect':
                answer = st.multiselect(
                    "Select all that apply:",
                    options=current_question['options'],
                    default=st.session_state.policy_answers.get(current_question['id'], []),
                    key=answer_key
                )
            
            # Store answer
            if answer:
                if current_question['type'] == 'multiselect':
                    st.session_state.policy_answers[current_question['id']] = answer
                else:
                    st.session_state.policy_answers[current_question['id']] = answer
            
            # Navigation buttons
            col_prev, col_next, col_cancel = st.columns([1, 3, 1])
            
            with col_prev:
                if st.button("◀ Previous", disabled=st.session_state.policy_question_index == 0, key=f"policy_prev_{answer_key}"):
                    st.session_state.policy_question_index -= 1
                    st.rerun()
            
            with col_next:
                is_disabled = False
                if current_question.get('required'):
                    if current_question['type'] == 'multiselect':
                        is_disabled = not answer or len(answer) == 0
                    else:
                        is_disabled = not answer or (isinstance(answer, str) and len(answer.strip()) == 0)
                
                if st.button("Next ▶", disabled=is_disabled, key=f"policy_next_{answer_key}"):
                    if is_disabled:
                        st.warning("Please answer this question to continue.")
                    else:
                        st.session_state.policy_question_index += 1
                        st.rerun()
            
            with col_cancel:
                if st.button("❌ Cancel", key=f"policy_cancel_{answer_key}"):
                    st.session_state.policy_questionnaire_active = False
                    st.session_state.policy_answers = {}
                    st.session_state.policy_question_index = 0
                    st.rerun()
        
        else:
            # All questions answered - generate policies
            st.success("✅ All questions answered!")
            
            st.subheader("📋 Your Answers Summary")
            for q in visible_questions:
                if q['id'] in st.session_state.policy_answers:
                    st.markdown(f"**{q['question']}**")
                    answer = st.session_state.policy_answers[q['id']]
                    if isinstance(answer, list):
                        st.markdown(f"  → {', '.join(answer)}")
                    else:
                        st.markdown(f"  → {answer}")
                    st.markdown("---")
            
            col_generate, col_restart = st.columns([3, 1])
            
            with col_generate:
                if st.button("📋 Generate Policies", type="primary"):
                    policy_agent, policy_error = load_policy_agent()
                    if policy_error:
                        st.error(f"Error: {policy_error}")
                    else:
                        with st.spinner("Generating policies..."):
                            # Prepare context from answers
                            policy_context = st.session_state.policy_answers.copy()
                            # Merge with company data if available
                            if st.session_state.company_data:
                                policy_context.update(st.session_state.company_data)
                            
                            # Determine which policies to generate
                            policy_priority = policy_context.get('policy_priority', [])
                            if isinstance(policy_priority, str):
                                policy_priority = [policy_priority]
                            
                            # Generate requested policies
                            responses = []
                            if 'Privacy Policy' in policy_priority:
                                response = policy_agent.generate_privacy_policy(policy_context)
                                responses.append(("Privacy Policy", response))
                            if 'Terms of Service' in policy_priority:
                                response = policy_agent.generate_terms_of_service(policy_context)
                                responses.append(("Terms of Service", response))
                            if not responses:
                                # Default to privacy policy
                                response = policy_agent.generate_privacy_policy(policy_context)
                                responses.append(("Privacy Policy", response))
                            
                            # Combine responses
                            combined_response = "# 📋 Generated Policies\n\n"
                            all_sources = []
                            for policy_name, resp in responses:
                                combined_response += f"## {policy_name}\n\n{resp['response']}\n\n---\n\n"
                                all_sources.extend(resp.get('sources', []))
                            
                            st.session_state.messages.append({
                                "role": "assistant",
                                "content": combined_response,
                                "sources": all_sources
                            })
                            
                            # Reset questionnaire
                            st.session_state.policy_questionnaire_active = False
                            st.session_state.policy_answers = {}
                            st.session_state.policy_question_index = 0
                            st.rerun()
            
            with col_restart:
                if st.button("🔄 Restart"):
                    st.session_state.policy_questionnaire_active = False
                    st.session_state.policy_answers = {}
                    st.session_state.policy_question_index = 0
                    st.rerun()
    
    # Team & Hiring Questionnaire Flow
    if "team_questionnaire_active" not in st.session_state:
        st.session_state.team_questionnaire_active = False
    if "team_answers" not in st.session_state:
        st.session_state.team_answers = {}
    if "team_question_index" not in st.session_state:
        st.session_state.team_question_index = 0
    
    if st.session_state.team_questionnaire_active:
        st.markdown("---")
        st.subheader("👥 Team & Hiring Questionnaire")
        
        questionnaire = get_team_questionnaire()
        
        # Filter questions based on dependencies
        visible_questions = [q for q in questionnaire if should_show_question(q, st.session_state.team_answers)]
        
        if st.session_state.team_question_index < len(visible_questions):
            current_question = visible_questions[st.session_state.team_question_index]
            
            st.progress((st.session_state.team_question_index + 1) / len(visible_questions))
            st.caption(f"Question {st.session_state.team_question_index + 1} of {len(visible_questions)}")
            
            st.markdown(f"### {current_question['question']}")
            
            answer_key = f"team_q_{current_question['id']}"
            answer = None
            
            # Render input based on type
            if current_question['type'] == 'textarea':
                answer = st.text_area(
                    "Your answer:",
                    value=st.session_state.team_answers.get(current_question['id'], ''),
                    height=150,
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'text':
                answer = st.text_input(
                    "Your answer:",
                    value=st.session_state.team_answers.get(current_question['id'], ''),
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'select':
                current_value = st.session_state.team_answers.get(current_question['id'])
                default_index = 0
                if current_value and current_value in current_question['options']:
                    default_index = current_question['options'].index(current_value)
                
                answer = st.selectbox(
                    "Select an option:",
                    options=current_question['options'],
                    index=default_index,
                    key=answer_key
                )
            elif current_question['type'] == 'multiselect':
                answer = st.multiselect(
                    "Select all that apply:",
                    options=current_question['options'],
                    default=st.session_state.team_answers.get(current_question['id'], []),
                    key=answer_key
                )
            
            # Store answer
            if answer:
                if current_question['type'] == 'multiselect':
                    st.session_state.team_answers[current_question['id']] = answer
                else:
                    st.session_state.team_answers[current_question['id']] = answer
            
            # Navigation buttons
            col_prev, col_next, col_cancel = st.columns([1, 3, 1])
            
            with col_prev:
                if st.button("◀ Previous", disabled=st.session_state.team_question_index == 0, key=f"team_prev_{answer_key}"):
                    st.session_state.team_question_index -= 1
                    st.rerun()
            
            with col_next:
                is_disabled = False
                if current_question.get('required'):
                    if current_question['type'] == 'multiselect':
                        is_disabled = not answer or len(answer) == 0
                    else:
                        is_disabled = not answer or (isinstance(answer, str) and len(answer.strip()) == 0)
                
                if st.button("Next ▶", disabled=is_disabled, key=f"team_next_{answer_key}"):
                    if is_disabled:
                        st.warning("Please answer this question to continue.")
                    else:
                        st.session_state.team_question_index += 1
                        st.rerun()
            
            with col_cancel:
                if st.button("❌ Cancel", key=f"team_cancel_{answer_key}"):
                    st.session_state.team_questionnaire_active = False
                    st.session_state.team_answers = {}
                    st.session_state.team_question_index = 0
                    st.rerun()
        
        else:
            # All questions answered - generate team analysis/job descriptions
            st.success("✅ All questions answered!")
            
            st.subheader("📋 Your Answers Summary")
            for q in visible_questions:
                if q['id'] in st.session_state.team_answers:
                    st.markdown(f"**{q['question']}**")
                    answer = st.session_state.team_answers[q['id']]
                    if isinstance(answer, list):
                        st.markdown(f"  → {', '.join(answer)}")
                    else:
                        st.markdown(f"  → {answer}")
                    st.markdown("---")
            
            col_analyze, col_restart = st.columns([3, 1])
            
            with col_analyze:
                if st.button("👥 Analyze Team Needs", type="primary"):
                    team_agent, team_error = load_team_agent()
                    if team_error:
                        st.error(f"Error: {team_error}")
                    else:
                        with st.spinner("Analyzing team needs and generating recommendations..."):
                            # Prepare context from answers
                            team_context = st.session_state.team_answers.copy()
                            # Merge with company data if available
                            if st.session_state.company_data:
                                team_context.update(st.session_state.company_data)
                            
                            # Perform team analysis
                            response = team_agent.analyze_team_needs(
                                company_context=team_context,
                                team_context=team_context
                            )
                            
                            st.session_state.messages.append({
                                "role": "assistant",
                                "content": response['response'],
                                "sources": response.get('sources', [])
                            })
                            
                            # Reset questionnaire
                            st.session_state.team_questionnaire_active = False
                            st.session_state.team_answers = {}
                            st.session_state.team_question_index = 0
                            st.rerun()
            
            with col_restart:
                if st.button("🔄 Restart"):
                    st.session_state.team_questionnaire_active = False
                    st.session_state.team_answers = {}
                    st.session_state.team_question_index = 0
                    st.rerun()
    
    # Patent Questionnaire Flow
    if "patent_questionnaire_active" not in st.session_state:
        st.session_state.patent_questionnaire_active = False
    if "patent_answers" not in st.session_state:
        st.session_state.patent_answers = {}
    if "patent_question_index" not in st.session_state:
        st.session_state.patent_question_index = 0
    
    if st.session_state.patent_questionnaire_active:
        st.markdown("---")
        st.subheader("🔬 Patent Analysis Questionnaire")
        
        questionnaire = get_patent_questionnaire()
        
        # Filter questions based on dependencies
        visible_questions = [q for q in questionnaire if should_show_question(q, st.session_state.patent_answers)]
        
        if st.session_state.patent_question_index < len(visible_questions):
            current_question = visible_questions[st.session_state.patent_question_index]
            
            st.progress((st.session_state.patent_question_index + 1) / len(visible_questions))
            st.caption(f"Question {st.session_state.patent_question_index + 1} of {len(visible_questions)}")
            
            st.markdown(f"### {current_question['question']}")
            
            answer_key = f"patent_q_{current_question['id']}"
            answer = None
            
            # Render input based on type
            if current_question['type'] == 'textarea':
                answer = st.text_area(
                    "Your answer:",
                    value=st.session_state.patent_answers.get(current_question['id'], ''),
                    height=150,
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'text':
                answer = st.text_input(
                    "Your answer:",
                    value=st.session_state.patent_answers.get(current_question['id'], ''),
                    key=answer_key,
                    placeholder=current_question.get('placeholder', '')
                )
            elif current_question['type'] == 'select':
                current_value = st.session_state.patent_answers.get(current_question['id'])
                default_index = 0
                if current_value and current_value in current_question['options']:
                    default_index = current_question['options'].index(current_value)
                
                answer = st.selectbox(
                    "Select an option:",
                    options=current_question['options'],
                    index=default_index,
                    key=answer_key
                )
            elif current_question['type'] == 'multiselect':
                answer = st.multiselect(
                    "Select all that apply:",
                    options=current_question['options'],
                    default=st.session_state.patent_answers.get(current_question['id'], []),
                    key=answer_key
                )
            
            # Store answer
            if answer:
                if current_question['type'] == 'multiselect':
                    st.session_state.patent_answers[current_question['id']] = answer
                else:
                    st.session_state.patent_answers[current_question['id']] = answer
            
            # Navigation buttons
            col_prev, col_next, col_cancel = st.columns([1, 3, 1])
            
            with col_prev:
                if st.button("◀ Previous", disabled=st.session_state.patent_question_index == 0, key=f"prev_{answer_key}"):
                    st.session_state.patent_question_index -= 1
                    st.rerun()
            
            with col_next:
                is_disabled = False
                if current_question.get('required'):
                    if current_question['type'] == 'multiselect':
                        is_disabled = not answer or len(answer) == 0
                    else:
                        is_disabled = not answer or (isinstance(answer, str) and len(answer.strip()) == 0)
                
                if st.button("Next ▶", disabled=is_disabled, key=f"next_{answer_key}"):
                    if is_disabled:
                        st.warning("Please answer this question to continue.")
                    else:
                        st.session_state.patent_question_index += 1
                        st.rerun()
            
            with col_cancel:
                if st.button("❌ Cancel", key=f"cancel_{answer_key}"):
                    st.session_state.patent_questionnaire_active = False
                    st.session_state.patent_answers = {}
                    st.session_state.patent_question_index = 0
                    st.rerun()
        
        else:
            # All questions answered - show summary and analyze
            st.success("✅ All questions answered!")
            
            st.subheader("📋 Your Answers Summary")
            for q in visible_questions:
                if q['id'] in st.session_state.patent_answers:
                    st.markdown(f"**{q['question']}**")
                    answer = st.session_state.patent_answers[q['id']]
                    if isinstance(answer, list):
                        st.markdown(f"  → {', '.join(answer)}")
                    else:
                        st.markdown(f"  → {answer}")
                    st.markdown("---")
            
            col_analyze, col_restart = st.columns([3, 1])
            
            with col_analyze:
                if st.button("🔬 Analyze Patent", type="primary"):
                    patent_agent, patent_error = load_patent_agent()
                    if patent_error:
                        st.error(f"Error: {patent_error}")
                    else:
                        with st.spinner("Analyzing your invention and generating comprehensive patent analysis..."):
                            # Prepare context from answers
                            invention_context = {
                                'invention_description': st.session_state.patent_answers.get('invention_description', ''),
                                'industry': st.session_state.patent_answers.get('industry', ''),
                                'novelty': st.session_state.patent_answers.get('novelty', ''),
                                'development_stage': st.session_state.patent_answers.get('development_stage', ''),
                                'prior_art_awareness': st.session_state.patent_answers.get('prior_art_awareness', ''),
                                'prior_art_details': st.session_state.patent_answers.get('prior_art_details', ''),
                                'public_disclosure': st.session_state.patent_answers.get('public_disclosure', ''),
                                'disclosure_date': st.session_state.patent_answers.get('disclosure_date', ''),
                                'filing_priority': st.session_state.patent_answers.get('filing_priority', ''),
                                'geographic_interest': st.session_state.patent_answers.get('geographic_interest', []),
                                'budget_range': st.session_state.patent_answers.get('budget_range', '')
                            }
                            
                            # Perform comprehensive analysis
                            invention_desc = f"{invention_context['invention_description']} - {invention_context['novelty']}"
                            
                            # 1. Assess patentability
                            patentability_response = patent_agent.assess_patentability(
                                invention_desc, 
                                invention_context
                            )
                            
                            # 2. Search patents
                            search_query = f"{invention_context['invention_description']} {invention_context['industry']}"
                            search_response = patent_agent.search_patents(
                                search_query,
                                invention_context
                            )
                            
                            # 3. Filing strategy
                            filing_response = patent_agent.filing_strategy(invention_context)
                            
                            # Combine all analyses
                            combined_response = f"""# 🔬 Comprehensive Patent Analysis

## 📋 Patentability Assessment
{patentability_response['response']}

---

## 🔍 Patent Search Results
{search_response['response']}

---

## 📄 Filing Strategy
{filing_response['response']}
"""
                            
                            # Combine sources
                            all_sources = (
                                patentability_response.get('sources', []) +
                                search_response.get('sources', []) +
                                filing_response.get('sources', [])
                            )
                            
                            st.session_state.messages.append({
                                "role": "assistant",
                                "content": combined_response,
                                "sources": all_sources
                            })
                            
                            # Save questionnaire answers as user choices
                            company_id = st.session_state.get("company_id", "default")
                            user_id = st.session_state.get("user_id", "anonymous")
                            
                            choice_manager.save_choice(
                                user_id=user_id,
                                company_id=company_id,
                                agent_name="patent",
                                action="questionnaire_answers",
                                choice_id="patent_questionnaire_completed",
                                choice_data={
                                    "label": "Patent Questionnaire Completed",
                                    "description": "User completed patent analysis questionnaire",
                                    "value": st.session_state.patent_answers.copy()
                                },
                                context={
                                    "questionnaire_type": "patent_analysis",
                                    "num_questions": len(visible_questions),
                                    "invention": invention_context.get('invention_description', '')[:100]
                                }
                            )
                            
                            # Reset questionnaire
                            st.session_state.patent_questionnaire_active = False
                            st.session_state.patent_answers = {}
                            st.session_state.patent_question_index = 0
                            st.rerun()
            
            with col_restart:
                if st.button("🔄 Restart"):
                    st.session_state.patent_questionnaire_active = False
                    st.session_state.patent_answers = {}
                    st.session_state.patent_question_index = 0
                    st.rerun()
    
    # Intelligent Pitch Questionnaire Flow
    if "pitch_questionnaire_active" not in st.session_state:
        st.session_state.pitch_questionnaire_active = False
    if "pitch_answers" not in st.session_state:
        st.session_state.pitch_answers = {}
    if "pitch_question_index" not in st.session_state:
        st.session_state.pitch_question_index = 0
    if "pitch_existing_data" not in st.session_state:
        st.session_state.pitch_existing_data = {}
    
    if st.session_state.pitch_questionnaire_active:
        st.markdown("---")
        st.subheader("🚀 Intelligent Pitch Builder - Let's Build Your Perfect Pitch!")
        
        # Show existing data summary if available
        if st.session_state.pitch_existing_data:
            with st.expander("✅ Information I Already Understand", expanded=False):
                st.info("Great! I already have some information about your startup. Let me ask targeted questions to fill in the gaps.")
                for q_id, data in list(st.session_state.pitch_existing_data.items())[:5]:  # Show first 5
                    st.markdown(f"**{data.get('question', '')}**")
                    value = data.get('value', '')
                    if isinstance(value, str) and len(value) > 100:
                        value = value[:100] + "..."
                    st.caption(f"✓ {value}")
        
        questions_to_ask = st.session_state.get("pitch_questions_to_ask", [])
        
        if st.session_state.pitch_question_index < len(questions_to_ask):
            current_question = questions_to_ask[st.session_state.pitch_question_index]
            analyzer = PitchQuestionAnalyzer()
            
            st.progress((st.session_state.pitch_question_index + 1) / len(questions_to_ask))
            st.caption(f"Question {st.session_state.pitch_question_index + 1} of {len(questions_to_ask)}")
            
            # Show contextual message
            contextual_msg = analyzer.get_contextual_message(current_question, st.session_state.pitch_existing_data)
            st.info(contextual_msg)
            
            st.markdown(f"### {current_question['question']}")
            
            answer_key = f"pitch_q_{current_question['id']}"
            answer = None
            
            # All pitch questions are textarea for detailed answers
            answer = st.text_area(
                "Your answer:",
                value=st.session_state.pitch_answers.get(current_question['id'], ''),
                height=150,
                key=answer_key,
                placeholder=f"Provide a detailed answer to help build your pitch..."
            )
            
            # Store answer
            if answer:
                st.session_state.pitch_answers[current_question['id']] = answer
            
            # Navigation buttons
            col_prev, col_next, col_cancel = st.columns([1, 3, 1])
            
            with col_prev:
                if st.button("◀ Previous", disabled=st.session_state.pitch_question_index == 0, key=f"pitch_prev_{answer_key}"):
                    st.session_state.pitch_question_index -= 1
                    st.rerun()
            
            with col_next:
                is_disabled = not answer or (isinstance(answer, str) and len(answer.strip()) == 0)
                
                if st.button("Next ▶", disabled=is_disabled, key=f"pitch_next_{answer_key}"):
                    if is_disabled:
                        st.warning("Please provide an answer to continue.")
                    else:
                        st.session_state.pitch_question_index += 1
                        st.rerun()
            
            with col_cancel:
                if st.button("❌ Cancel", key=f"pitch_cancel_{answer_key}"):
                    st.session_state.pitch_questionnaire_active = False
                    st.session_state.pitch_answers = {}
                    st.session_state.pitch_question_index = 0
                    st.session_state.pitch_existing_data = {}
                    st.rerun()
        
        else:
            # All questions answered - merge answers and generate pitch options
            st.success("✅ All questions answered! Now let's create your pitch options.")
            
            # Merge answers with company data
            analyzer = PitchQuestionAnalyzer()
            merged_company_data = analyzer.merge_answers_with_company_data(
                st.session_state.company_data,
                st.session_state.pitch_answers
            )
            st.session_state.company_data = merged_company_data
            
            # Generate 2-3 pitch options
            st.subheader("🎯 Generating Pitch Options...")
            with st.spinner("Creating 3 different pitch deck options with different angles..."):
                try:
                    pitch_options = generate_pitch_options(agent, merged_company_data, num_options=3)
                    
                    if pitch_options:
                        st.session_state.pitch_options = pitch_options
                        st.session_state.pitch_questionnaire_active = False  # Close questionnaire
                        st.session_state.pitch_options_ready = True  # Show options selection
                        st.rerun()
                    else:
                        st.error("Failed to generate pitch options. Please try again.")
                except Exception as e:
                    st.error(f"Error generating pitch options: {e}")
                    logger.error(f"Error in pitch options generation: {e}", exc_info=True)
    
    # Pitch Options Selection and Refinement Flow
    if st.session_state.get("pitch_options_ready") and st.session_state.get("pitch_options"):
        st.markdown("---")
        st.subheader("🎯 Choose Your Pitch Style")
        st.info("I've created 3 different pitch deck options. Review and select the one that best fits your needs, then we'll refine it together!")
        
        pitch_options = st.session_state.pitch_options
        selected_option_index = st.session_state.get("selected_pitch_option_index", None)
        
        # Display options in columns
        cols = st.columns(len(pitch_options))
        
        for idx, option in enumerate(pitch_options):
            with cols[idx]:
                is_selected = selected_option_index == idx
                st.markdown(f"### {option['title']}")
                st.caption(option['description'])
                st.info(f"📊 {option['slides_count']} slides")
                
                # Show preview of first slide
                if option.get('slides_data', {}).get('slides'):
                    first_slide = option['slides_data']['slides'][0]
                    with st.expander("Preview First Slide", expanded=False):
                        st.markdown(f"**{first_slide.get('title', '')}**")
                        st.caption(first_slide.get('content', '')[:200] + "...")
                
                if st.button(f"Select Option {idx + 1}", key=f"select_option_{idx}", 
                           type="primary" if is_selected else "secondary"):
                    st.session_state.selected_pitch_option_index = idx
                    st.session_state.selected_pitch_option = option
                    st.rerun()
        
        # If option selected, show refinement step
        if selected_option_index is not None:
            st.markdown("---")
            st.subheader("✨ Refine Your Selected Pitch")
            selected_option = st.session_state.selected_pitch_option
            
            st.success(f"✅ You selected: **{selected_option['title']}**")
            st.info("💡 Would you like to make any changes to this pitch? Tell me what you'd like to adjust, and I'll update it.")
            
            refinement_request = st.text_area(
                "What would you like to change? (e.g., 'Emphasize our traction more', 'Add more details about the market opportunity', 'Make the problem statement more urgent')",
                height=100,
                key="pitch_refinement_request",
                placeholder="Describe the changes you'd like..."
            )
            
            col_refine, col_finalize = st.columns(2)
            
            with col_refine:
                if st.button("🔧 Apply Refinements", disabled=not refinement_request or len(refinement_request.strip()) < 10):
                    if refinement_request:
                        with st.spinner("Applying your refinements..."):
                            # Get selected slides
                            selected_slides = selected_option['slides_data']['slides']
                            
                            # Use agent to refine based on request
                            # This is a simplified version - you might want to create a refine_slides method
                            try:
                                # Merge refinement request into company data
                                refined_company_data = {
                                    **st.session_state.company_data,
                                    "refinement_request": refinement_request,
                                    "existing_slides": selected_slides,
                                    "pitch_angle": selected_option['style']
                                }
                                
                                # Generate refined slides
                                refined_slides_data = agent.generate_slides(refined_company_data, gamma_only=True)
                                
                                # Update selected option
                                st.session_state.selected_pitch_option['slides_data'] = refined_slides_data
                                st.session_state.selected_pitch_option['slides_count'] = refined_slides_data.get('total_slides', 0)
                                st.session_state.pitch_refinement_history = st.session_state.get('pitch_refinement_history', [])
                                st.session_state.pitch_refinement_history.append(refinement_request)
                                
                                st.success("✅ Refinements applied! Review the updated pitch below.")
                                st.rerun()
                            except Exception as e:
                                st.error(f"Error applying refinements: {e}")
                                logger.error(f"Error refining pitch: {e}", exc_info=True)
            
            with col_finalize:
                if st.button("✅ Finalize & Generate", type="primary"):
                    # User is happy with the pitch, now generate final version with PPT and Gamma
                    with st.spinner("Generating final pitch deck with PowerPoint and Gamma presentation..."):
                        selected_slides = st.session_state.selected_pitch_option['slides_data']['slides']
                        company_data_final = {
                            **st.session_state.company_data,
                            "enhance_with_ai": st.session_state.get('enhance_with_ai', True),
                            "ai_provider": st.session_state.get('ai_provider', 'Auto'),
                            "gamma_theme": st.session_state.get("gamma_theme", "startup-pitch"),
                            "final_slides": selected_slides  # Use refined slides
                        }
                        
                        # Generate final slides (with PPT and Gamma)
                        try:
                            # Create slides data structure
                            final_slides_data = {
                                "slides": selected_slides,
                                "total_slides": len(selected_slides),
                                "company_name": st.session_state.company_data.get('company_name', 'Unknown')
                            }
                            
                            # Generate PowerPoint
                            from src.utils.exporters import PitchExporter
                            exporter = PitchExporter()
                            enhance_with_ai = st.session_state.get('enhance_with_ai', True)
                            ai_provider = st.session_state.get('ai_provider', 'Auto')
                            
                            pptx_path = exporter.export_to_powerpoint(
                                final_slides_data["slides"],
                                final_slides_data["company_name"],
                                include_images=True,
                                enhance_with_ai=enhance_with_ai,
                                company_data=st.session_state.company_data,
                                full_rewrite=True,
                                ai_provider=ai_provider
                            )
                            
                            if pptx_path:
                                final_slides_data["pptx_path"] = pptx_path
                            
                            # Generate Gamma (use existing slides data)
                            from src.utils.gamma_integration import GammaIntegration
                            gamma = GammaIntegration()
                            
                            gamma_result = gamma.create_presentation(
                                final_slides_data["slides"],
                                final_slides_data["company_name"],
                                theme_id=st.session_state.get("gamma_theme", "startup-pitch"),
                                enhance_with_ai=enhance_with_ai
                            )
                            
                            final_slides_data["gamma_presentation"] = gamma_result
                            
                            # Store final slides
                            st.session_state.slides = final_slides_data
                            st.session_state.current_slide = 0
                            
                            # Reset pitch questionnaire state
                            st.session_state.pitch_questionnaire_active = False
                            st.session_state.pitch_options_ready = False
                            st.session_state.selected_pitch_option_index = None
                            st.session_state.pitch_options = None
                            st.session_state.selected_pitch_option = None
                            
                            st.success("🎉 **Pitch deck generated successfully!** Check the slides viewer and download options.")
                            st.rerun()
                            
                        except Exception as e:
                            st.error(f"Error generating final pitch: {e}")
                            logger.error(f"Error in final pitch generation: {e}", exc_info=True)
    
    # Additional actions
    col5, col6 = st.columns(2)
    with col5:
        if st.button("💡 Pitch Tips"):
            response = agent.process_query("What are common pitch deck mistakes and best practices?")
            st.session_state.messages.append({
                "role": "assistant",
                "content": response['response'],
                "sources": response.get('sources', [])
            })
            st.rerun()
    
    with col6:
        if st.button("📈 Enhanced Scoring"):
            if st.session_state.messages:
                # Get last pitch from messages
                last_pitch = None
                for msg in reversed(st.session_state.messages):
                    if msg.get("role") == "assistant" and len(msg.get("content", "")) > 100:
                        last_pitch = msg.get("content", "")
                        break
                
                if last_pitch:
                    with st.spinner("Scoring pitch..."):
                        response = agent.evaluate_pitch_with_scores(
                            last_pitch,
                            st.session_state.company_data
                        )
                        eval_data = response.get('evaluation', {})
                        
                        # Display scores
                        st.subheader("📊 Pitch Scores")
                        overall = eval_data.get('overall_score', 0)
                        st.metric("Overall Score", f"{overall}/10")
                        
                        section_scores = eval_data.get('section_scores', {})
                        for section, score_data in section_scores.items():
                            if isinstance(score_data, dict):
                                score = score_data.get('score', 0)
                                st.progress(score / 10)
                                st.caption(f"{section.title()}: {score}/10 - {score_data.get('feedback', '')}")
                        
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": f"## Enhanced Scoring\n\n{json.dumps(eval_data, indent=2)}",
                            "sources": response.get('sources', [])
                        })
                        st.rerun()
                else:
                    st.info("Generate a pitch first to score it!")


if __name__ == "__main__":
    main()

